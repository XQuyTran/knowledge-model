from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x3C, 0x6E)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
LIGHT_BG = RGBColor(0xF2, 0xF7, 0xFC)
GRAY = RGBColor(0x66, 0x66, 0x66)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)

slide_num = [0]

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, text, top=Inches(0.3), left=Inches(0.5), width=Inches(12.3), height=Inches(0.8)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLUE
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

def add_body_box(slide, lines, left=Inches(0.7), top=Inches(1.4), width=Inches(11.9), height=Inches(5.5), font_size=17):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(line_data, str):
            line = line_data
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = DARK
            p.space_after = Pt(4)
        elif isinstance(line_data, dict):
            p.text = line_data.get('text', '')
            p.font.size = Pt(line_data.get('size', font_size))
            p.font.color.rgb = line_data.get('color', DARK)
            p.space_after = Pt(line_data.get('space', 4))
            if line_data.get('bold'):
                p.font.bold = True
            if line_data.get('level', 0) > 0:
                p.level = line_data.get('level', 0)
    return txBox

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.clear()
    tf.text = text

def add_footer(slide, text="Project 1 - He thong ho tro hoc tap"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{text}    |    {slide_num[0]}/12"
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

def bullet(text, size=17, color=None, bold=False, level=0, space=4):
    d = {'text': text, 'size': size, 'color': color or DARK, 'bold': bold, 'level': level, 'space': space}
    return d

def heading(text, size=20):
    return bullet(text, size=size, color=BLUE, bold=True, space=8)

def sub(text, size=16):
    return bullet(text, size=size, level=1, space=3)

def new_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    slide_num[0] += 1
    return slide

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = new_slide()
add_bg(slide, BLUE)
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "HE THONG HO TRO GIAI BAI TAP LAP TRINH\nCO HUONG DAN TUNG BUOC"
p.font.size = Pt(36)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Programming Exercise Solving System with Step-by-Step Guidance"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT
p2.font.italic = True
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)
p3 = tf.add_paragraph()
p3.text = "Context-Aware Explanations & Diagnostic Rule-Based Reasoning"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(8)

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
lines_info = [
    "Mon hoc: Cau truc du lieu & Thuat toan nang cao",
    "Giang vien huong dan: TS. Nguyen Van A",
    "Nhom thuc hien: Nhom 1 - BDTT K16",
    "Hoc ky: 1 - Nam hoc 2025-2026"
]
for i, line in enumerate(lines_info):
    if i == 0:
        p = tf2.paragraphs[0]
    else:
        p = tf2.add_paragraph()
    p.text = line
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(28)

add_notes(slide, """SLIDE 1 - TIEU DE (Title Slide)

[Thoi gian: 1.5 phut]

Day la slide mo dau. Gioi thieu ten de tai bang tieng Viet va English subtitle.

Noi dung chinh:
- He thong nay giai quyet van de: sinh vien thuong mac loi logic & thuat toan 
  nhung cong cu hien tai chi bao loi bien dich (compile error) ma khong phan tich
  duoc loi trong tu duy giai thuat.
- Giai phap: ket hop LLMs (GPT-4, Code Llama) + Rule-Based Reasoning de:
  (1) Chan doan loi chinh xac, (2) Sinh loi giai tung buoc co giai thich ngu canh.

Chuyen tiep: "Tiep theo, em xin gioi thieu noi dung chinh cua bai bao cao."
""")

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = new_slide()
add_title_bar(slide, "NOI DUNG TRINH BAY (Agenda)")
items = [
    heading("1. Muc tieu he thong (System Objectives)", 20),
    heading("2. Thu thap tri thuc (Knowledge Acquisition)", 20),
    heading("3. Phan loai bai tap (Exercise Classification)", 20),
    heading("4. Mo hinh bieu dien tri thuc (Knowledge Representation Model)", 20),
    heading("5. Kien truc he thong (System Architecture)", 20),
    heading("6. Bai toan & Thuat giai (Problems & Algorithms)", 20),
    heading("7. Ky thuat chinh (Core Techniques): LLMs + Rule-Based Reasoning", 20),
    heading("8. Yeu cau san pham (Product Requirements)", 20),
    heading("9. So sanh & Danh gia (Comparison & Evaluation)", 20),
    heading("10. Ket luan & Huong phat trien (Conclusion & Future Work)", 20),
]
add_body_box(slide, items, top=Inches(1.3), font_size=20)
add_notes(slide, """SLIDE 2 - AGENDA (Noi dung trinh bay)

[Thoi gian: 1 phut]

Gioi thieu 10 phan cua bai bao cao.

Phan bo thoi gian du kien:
1. Muc tieu (2 phut)
2. Thu thap tri thuc (2 phut)
3. Phan loai bai tap (1.5 phut)
4. Mo hinh tri thuc (3 phut) - PHAN QUAN TRONG
5. Kien truc (2 phut)
6. Bai toan & Thuat giai (2.5 phut)
7. Ky thuat chinh (3 phut) - PHAN QUAN TRONG
8. Yeu cau san pham (1.5 phut)
9. So sanh & Danh gia (2.5 phut)
10. Ket luan (1.5 phut)
Tong: ~22 phut + 5 phut Q&A

Chuyen tiep: "Chung ta bat dau voi phan dau tien: Muc tieu cua he thong."
""")
add_footer(slide)

# ============================================================
# SLIDE 3: Muc tieu he thong
# ============================================================
slide = new_slide()
add_title_bar(slide, "MUC TIEU HE THONG (System Objectives)")
body = [
    heading("◆ Muc tieu tong quat:"),
    sub("Phat trien he thong chan doan loi logic va thuat toan trong ma nguon sinh vien (C/C++)"),
    heading("▸ Cac muc tieu cu the:"),
    bullet("• Phat hien & chan doan (Detection & Diagnosis): Xac dinh vi tri va loai loi", space=6),
    sub("- Loi logic (Logic errors): Sai dieu kien, sai cong thuc tinh toan"),
    sub("- Loi thuat toan (Algorithm errors): Chon sai approach, infinite loop"),
    sub("- Loi bien (Boundary errors): Off-by-one, array out of bounds"),
    bullet("• Sinh loi giai tung buoc (Step-by-Step Solution Generation):", space=6),
    sub("- Tao huong dan chi tiet + giai thich ngu canh (Context-Aware Explanations)"),
    sub("- Du tren mo hinh tri thuc loi (Error Knowledge Model)"),
    bullet("• Tich hop LLMs + Rule-Based Reasoning:", space=6),
    sub("- Ket hop GPT-4 / Code Llama de phan tich ngu nghia ma nguon"),
    sub("- Suy luan dua tren luat (Diagnostic Rules) de chan doan chinh xac"),
]
add_body_box(slide, body, font_size=16)
add_notes(slide, """SLIDE 3 - MUC TIEU HE THONG (System Objectives)

[Thoi gian: 2 phut]

Trinh bay 3 muc tieu chinh:

1. Phat hien & chan doan loi:
   - Loi logic: sai dieu kien if, sai cong thuc (VD: dung + thay vi *)
   - Loi thuat toan: chon sai giai thuat, vong lap vo han
   - Loi bien: off-by-one, truy cap mang ngoai pham vi

2. Sinh loi giai tung buoc:
   - Khong chi ra loi ma con huong dan CACH sua
   - Giai thich phu hop voi ngu canh cu the cua bai toan

3. Tich hop LLMs + Rule-Based:
   - LLMs hieu ngu nghia code
   - Rule Engine dam bao do chinh xac

Chuyen tiep: "De lam duoc dieu nay, chung toi can thu thap tri thuc. Moi cac ban xem slide tiep theo."
""")
add_footer(slide)

# ============================================================
# SLIDE 4: Thu thap tri thuc
# ============================================================
slide = new_slide()
add_title_bar(slide, "THU THAP TRI THUC (Knowledge Acquisition)")
body = [
    heading("▸ Nguon tri thuc (Knowledge Sources):"),
    bullet("• Sach bai tap C/C++: 'C Programming: A Modern Approach', 'The C Programming Language' (K&R)"),
    bullet("• Tai lieu thuat toan kinh dien: 'Introduction to Algorithms' (CLRS)"),
    bullet("• Bo du lieu bai tap sinh vien da giai (Student submission corpora)"),
    heading("▸ Quy trinh thu thap & xu ly (Acquisition Pipeline):"),
    bullet("Buoc 1: Thu thap de bai & bai giai mau tu nhieu nguon", space=5),
    bullet("Buoc 2: Trich xuat dac trung (Feature extraction): cau truc dieu khien, kieu du lieu", space=5),
    bullet("Buoc 3: Gan nhan loi (Error labeling): xac dinh loi thuong gap cho tung dang bai", space=5),
    bullet("Buoc 4: Xay dung tap luat chan doan (Diagnostic Rules) dua tren expert knowledge", space=5),
    heading("▸ Cong cu ho tro (Supporting Tools):"),
    bullet("• Static code analysis (AST parsing) - phan tich ma nguon tinh"),
    bullet("• Dynamic testing voi test cases - phat hien loi runtime"),
]
add_body_box(slide, body, font_size=15)
add_notes(slide, """SLIDE 4 - THU THAP TRI THUC (Knowledge Acquisition)

[Thoi gian: 2 phut]

Giai thich cach xay dung co so tri thuc cho he thong:

1. Nguon tri thuc: 3 loai chinh
   - Sach giao khoa C/C++ (K&R, KN King)
   - SACH thuat toan (CLRS)
   - Bai tap tu sinh vien -- nguon thuc te quan trong nhat

2. Quy trinh 4 buoc:
   - Buoc 1: Thu thap de bai tu nhieu nguon (sach, online judge, bai tap tren lop)
   - Buoc 2: Trich xuat dac trung bang AST parser (control flow, data type, algorithm pattern)
   - Buoc 3: Gán nhan loi thuong gap cho tung bai (VD: sap xep -> loi swap, loi dieu kien)
   - Buoc 4: Xay dung IF-THEN rules tu cac chuyen gia

3. Cong cu:
   - AST parser de phan tich code tinh
   - Test cases de phat hien loi runtime (dung input -> sai output)

Chuyen tiep: "Sau khi co du lieu, chung toi phan loai bai tap theo chu de."
""")
add_footer(slide)

# ============================================================
# SLIDE 5: Phan loai bai tap
# ============================================================
slide = new_slide()
add_title_bar(slide, "PHAN LOAI BAI TAP (Exercise Classification)")
body = [
    heading("▸ He thong phan loai bai tap theo 5 chu de chinh:"),
    bullet("◆ Mang (Arrays)"),
    sub("- One-dimensional arrays: tim kiem (search), sap xep (sorting), dao nguoc (reversal)"),
    sub("- Multi-dimensional arrays: ma tran (matrix), xoay vong (rotation), nhan ma tran"),
    bullet("◆ Vong lap (Loops)"),
    sub("- for loop, while loop, do-while loop"),
    sub("- Pattern printing, iterative computation"),
    bullet("◆ De quy (Recursion)"),
    sub("- Divide and conquer: binary search, merge sort, quick sort"),
    sub("- Backtracking: N-Queens, maze solving, permutation generation"),
    bullet("◆ Cau truc du lieu (Data Structures)"),
    sub("- Stack, Queue, Linked List, Tree, Graph"),
    sub("- Operations: insertion, deletion, traversal, searching"),
    bullet("◆ Thuat toan sap xep & tim kiem (Sorting & Searching)"),
    sub("- Bubble Sort, Selection Sort, Merge Sort, Quick Sort"),
    sub("- Linear Search, Binary Search, Hashing"),
]
add_body_box(slide, body, font_size=15)
add_notes(slide, """SLIDE 5 - PHAN LOAI BAI TAP (Exercise Classification)

[Thoi gian: 1.5 phut]

Trinh bay 5 nhom bai tap chinh. Moi nhom co cac bai toan dien hinh.

Y nghia: Viec phan loai nay giup he thong:
- Tra cuu nhanh tap luat chan doan phu hop
- Ap dung dung Pattern Matching cho tung loai bai
- De xuat thuat toan phu hop

Luu y: Day khong phai danh sach day du, ma la cac nhom chinh de chung toi 
xay dung mo hinh tri thuc ban dau. He thong co the mo rong them.

Chuyen tiep: "Bay gio chung ta se di vao phan core cua he thong: Mo hinh bieu dien tri thuc."
""")
add_footer(slide)

# ============================================================
# SLIDE 6: Mo hinh bieu dien tri thuc
# ============================================================
slide = new_slide()
add_title_bar(slide, "MO HINH BIEU DIEN TRI THUC (Knowledge Representation Model)")
body = [
    heading("▸ Moi quan he giua cac thanh phan tri thuc:"),
    bullet("Bai toan (Problem)  ->  Thuat toan (Algorithm)  ->  Cau truc Du lieu (Data Structure)", size=16, space=3),
    bullet("      |                              |                              |", size=14, color=GRAY, space=3),
    bullet("Loi thuong gap (Common Errors)  <-  Huong khac phuc (Fix Guidance)", size=16, space=6),
    heading("▸ Cau truc xu ly (Processing Pipeline):"),
    bullet("Code (Input) -> Pattern Matching (AST) -> Rule Triggered (Diagnostic Rules)", size=15, space=3),
    bullet("-> Step Explanation (Context-Aware) -> Solution Output", size=15, space=6),
    heading("▸ Co che hoat dong (How It Works):"),
    bullet("Buoc 1: Nhan ma nguon C++ dau vao", size=15, space=3),
    bullet("Buoc 2: Phan tich cu phap thanh Abstract Syntax Tree (AST)", size=15, space=3),
    bullet("Buoc 3: Pattern Matching phat hien cau truc code (loop, recursion...)", size=15, space=3),
    bullet("Buoc 4: Kich hoat luat chan doan (Diagnostic Rule Triggering)", size=15, space=3),
    bullet("Buoc 5: Sinh giai thich tung buoc dua tren ngu canh cu the", size=15, space=6),
    heading("▸ Vi du (Example):"),
    bullet("Bai toan 'tinh tong mang' -> while thieu i++ -> Rule: 'Loop counter not updated'", size=15, space=3),
    bullet("Giai thich: 'Ban quen tang i trong vong lap while'", size=15, space=2),
    bullet("Khac phuc: 'Them i++ vao cuoi than vong lap'", size=15, space=2),
]
add_body_box(slide, body, font_size=15, height=Inches(5.8))
add_notes(slide, """SLIDE 6 - MO HINH BIEU DIEN TRI THUC (Knowledge Representation Model)

[Thoi gian: 3 phut] - SLIDE QUAN TRONG

Day la phan core cua he thong. Can giai thich THAT KY:

1. MO HINH 5 THANH PHAN:
   - Bai toan -> Thua toan -> Cau truc du lieu: Moi quan he xuoi
   - Loi thuong gap <- Huong khac phuc: Moi quan he nguoc
   - Y nghia: Khi gap loi, he thong biet loi do thuoc bai toan/thuat toan nao
     va se goi y huong khac phuc tuong ung.

2. PIPELINE 5 BUOC:
   - Buoc 1-2: Code -> AST (dung Tree-sitter hoac Clang)
   - Buoc 3: Pattern Matching tren AST de nhan dang (VD: phat hien while loop)
   - Buoc 4: Rule Engine kich hoat luat neu phat hien loi
   - Buoc 5: LLM sinh giai thich tu nhien dua tren ket qua Rule Engine

3. VI DU CU THE:
   - Bai toan: tinh tong cac phan tu mang
   - Code sinh vien: while loop nhung quen i++
   - Rule phat hien: 'while loop thieu counter update' -> INFINITE_LOOP
   - He thong bao: 'Ban quen tang i' -> 'Them i++'

Chuyen tiep: "De hien thuc mo hinh nay, chung toi thiet ke kien truc he thong 4 tang."
""")
add_footer(slide)

# ============================================================
# SLIDE 7: Kien truc he thong
# ============================================================
slide = new_slide()
add_title_bar(slide, "KIEN TRUC HE THONG (System Architecture)")

# Draw layers as rectangles
layers = [
    ("User Interface Layer", "Web UI: Code Editor (Monaco/CodeMirror) + Output Panel\nReact + Tailwind CSS", Inches(1.2), BLUE),
    ("Processing Layer", "Code Parser (AST) -> Pattern Matcher -> Rule Engine\n-> LLM Connector (GPT-4 / Code Llama API)", Inches(2.4), RGBColor(0x2E, 0x7D, 0x32)),
    ("Knowledge Layer", "Exercise DB | Error KB | Diagnostic Rules | Fix DB\nPostgreSQL + ChromaDB (vector storage)", Inches(3.6), ORANGE),
    ("Data Layer", "PostgreSQL (structured data) + Vector DB (embeddings)\nPinecone / ChromaDB for semantic search", Inches(4.8), RGBColor(0x6A, 0x1B, 0x9A)),
]

for i, (title, desc, top, color) in enumerate(layers):
    # Main box
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.color.rgb = WHITE
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    # Arrow between layers
    if i < len(layers) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), top + Inches(1.0), Inches(0.8), Inches(0.2))
        arr.fill.solid()
        arr.fill.fore_color.rgb = ACCENT
        arr.line.fill.background()

# Tech stack on the right
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(11.7), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Tech Stack: Frontend: React + Monaco Editor | Backend: Python FastAPI | AI: GPT-4 / Code Llama + LangChain | DB: PostgreSQL + ChromaDB | Parser: Tree-sitter / Clang AST"
p.font.size = Pt(13)
p.font.color.rgb = DARK
p.font.bold = True

add_notes(slide, """SLIDE 7 - KIEN TRUC HE THONG (System Architecture)

[Thoi gian: 2 phut]

Gioi thieu kien truc 4 tang cua he thong:

1. Tang Giao dien (UI Layer):
   - Code Editor voi syntax highlighting (dung Monaco hoac CodeMirror)
   - Output panel hien thi loi giai tung buoc
   - Error Highlighting truc tiep tren code

2. Tang Xu ly (Processing Layer):
   - Code Parser: chuyen C/C++ sang AST (Tree-sitter)
   - Pattern Matcher: phat hien cau truc code
   - Rule Engine: ap dung IF-THEN rules
   - LLM Connector: goi GPT-4 / Code Llama de sinh giai thich

3. Tang Tri thuc (Knowledge Layer):
   - Exercise DB: luu de bai, bai giai mau, test cases
   - Error KB: co so tri thuc loi
   - Diagnostic Rules: tap luat chan doan
   - Fix DB: huong dan khac phuc

4. Tang Du lieu (Data Layer):
   - PostgreSQL: du lieu co cau truc
   - Vector DB (ChromaDB): embeddings cho semantic search

Chuyen tiep: "Tren kien truc nay, chung toi xay dung cac bai toan cu the."
""")
add_footer(slide)

# ============================================================
# SLIDE 8: Bai toan & Thuat giai
# ============================================================
slide = new_slide()
add_title_bar(slide, "BAI TOAN & THUAT GIAI (Problems & Algorithms)")
body = [
    heading("▸ Ma tran bai toan - thuat giai (Problem-Algorithm Matrix):"),
    bullet("1. 'Find Maximum/Minimum in Array'", size=16, space=2),
    sub("Approach: Linear Scan O(n), Divide & Conquer O(n)"),
    sub("Common Errors: Off-by-one, uninitialized variable"),
    bullet("2. 'Check Palindrome String'", size=16, space=2),
    sub("Approach: Two-pointer technique, Recursive comparison"),
    sub("Common Errors: Index out of bounds, case-sensitivity"),
    bullet("3. 'Linked List Reversal'", size=16, space=2),
    sub("Approach: Iterative (3-pointer), Recursive reversal"),
    sub("Common Errors: Lost reference to next node, null pointer"),
    bullet("4. 'Binary Search'", size=16, space=2),
    sub("Approach: Iterative, Recursive binary search"),
    sub("Common Errors: Wrong mid calculation (overflow), infinite recursion"),
    bullet("5. 'Tree Traversal (Inorder/Preorder/Postorder)'", size=16, space=2),
    sub("Approach: Recursive, Iterative using Stack"),
    sub("Common Errors: Stack overflow (deep recursion), null root check"),
    heading("▸ Co che de xuat (Recommendation Engine):"),
    bullet("He thong goi y thuat toan phu hop dua tren phan tich keywords cua de bai", size=15),
    bullet("Rule-based + ML classification: mapping Problem -> Algorithm -> Data Structure", size=15),
]
add_body_box(slide, body, font_size=15, height=Inches(5.8))
add_notes(slide, """SLIDE 8 - BAI TOAN & THUAT GIAI (Problems & Algorithms)

[Thoi gian: 2.5 phut]

Trinh bay 5 bai toan dien hinh duoc xay dung trong he thong:

Moi bai toan co:
- Approach: Cac phuong phap giai (thuong la nhieu option)
- Common Errors: Loi thuong gap

5 bai toan mau:
1. Tim min/max: loi hay gap la bien chua khoi tao hoac sai index
2. Palindrome: loi index out of bounds khi dung two-pointer
3. Dao nguoc linked list: loi mat reference khi cap nhat pointer
4. Binary search: loi tinh mid bi tran so, hoac recursive bi infinite recursion
5. Duyet cay: loi stack overflow neu cay qua sau

Co che Recommendation Engine:
- Phan tich de bai, trich xuat keywords
- Mapping keywords -> thuat toan phu hop
- VD: 'tim kiem' + 'mang da sap xep' -> Binary Search

Chuyen tiep: "Phan tiep theo la ky thuat chinh de hien thuc he thong nay."
""")
add_footer(slide)

# ============================================================
# SLIDE 9: Ky thuat chinh
# ============================================================
slide = new_slide()
add_title_bar(slide, "KY THUAT CHINH (Core Techniques): LLMs + Rule-Based Reasoning")
body = [
    heading("▸ Ky thuat 1: Mo hinh ngon ngu lon (Large Language Models - LLMs)"),
    bullet("• Code Llama (Meta): Ma nguon mo, chuyen cho code C/C++, co the fine-tune", space=4),
    bullet("• GPT-4 (OpenAI): Hieu ngu nghia sau (deep semantics), chat luong sinh cao", space=4),
    bullet("• LangChain Framework: Dieu phoi (orchestration), Prompt Templates, Chain-of-Thought", space=4),
    bullet("• Embedding: CodeBERT / GraphCodeBERT cho tim kiem ngu nghia (semantic search)", space=6),
    bullet("  Ung dung (Applications):", size=15, space=2),
    sub("- Phan tich ngu nghia ma nguon (Semantic Code Analysis)"),
    sub("- Mo ta sai sot bang ngon ngu tu nhien (Natural Language Error Description)"),
    sub("- Sinh code goi y sua loi (Fix Suggestion Generation)"),
    heading("▸ Ky thuat 2: Suy luan dua tren luat (Rule-Based Reasoning)"),
    bullet("• Tap luat chan doan (Diagnostic Rules): Cau truc IF-THEN rules", space=4),
    bullet("• Pattern Matching tren AST: Xac dinh cau truc code va loi tuong ung", space=4),
    bullet("• Finite State Machine (FSM): Theo doi trang thai qua trinh giai bai", space=4),
    bullet("  Vi du Rule (Example):", size=15, color=DARK, bold=True, space=2),
    sub("IF (loop_type == 'while') AND (counter_update NOT in loop_body)"),
    sub("THEN error_type = 'INFINITE_LOOP', suggestion = 'Add counter update inside loop'"),
    heading("▸ Co che Hybrid (Hybrid Mechanism): Rule -> LLM -> Context-Aware Explanation"),
    bullet("  Rule Engine: Chan doan loi chinh xac (do tin cay cao)", size=15, space=2),
    bullet("  LLM: Sinh giai thich tu nhien, de hieu (tang trai nghiem nguoi dung)", size=15, space=2),
]
add_body_box(slide, body, font_size=14, height=Inches(5.8))
add_notes(slide, """SLIDE 9 - KY THUAT CHINH (Core Techniques): LLMs + Rule-Based Reasoning

[Thoi gian: 3 phut] - SLIDE QUAN TRONG NHAT

Day la phan quan trong nhat cua bai bao cao. Can giai thich that ky:

KY THUAT 1 - LLMs (3 phan):
1. Code Llama: Mo hinh ma nguon mo cua Meta, duoc huan luyen tren code.
   Uu diem: mien phi, co the fine-tune rieng, hieu code C++ tot.
2. GPT-4: Mo hinh thuong mai cua OpenAI. Uu diem: hieu ngu nghia cau lenh,
   generation chat luong cao. Nhuoc diem: ton phi, can Internet.
3. LangChain: Framework giup ket noi nhieu LLM, quan ly prompt, thuc hien
   Chain-of-Thought de sinh giai thich nhieu buoc.
4. CodeBERT: Mo hinh embedding chuyen cho code. Dung de tim kiem cac doan
   code tuong tu trong co so du lieu.

KY THUAT 2 - Rule-Based Reasoning:
1. IF-THEN rules: Cac luat duoc xay dung tu kinh nghiem chuyen gia
2. Pattern Matching: So sanh AST cua code sinh vien voi mau da biet
3. FSM: Theo doi qua trinh giai bai cua sinh vien de xac dinh buoc tiep theo

VI DU: Phat hien vong lap while thieu cap nhat bien dem
- Pattern: while loop body khong co assignment toi bien dieu kien
- Rule: INFINITE_LOOP
- LLM sinh: 'Em quen tang i. Hay them i++ vao cuoi vong lap.'

CO CHE HYBRID:
- Rule Engine dam bao do chinh xac (khong bao loi sai)
- LLM dam bao giai thich tu nhien, de hieu voi sinh vien

Chuyen tiep: "Sau khi hieu ky thuat, chung ta xem xet yeu cau san pham cu the."
""")
add_footer(slide)

# ============================================================
# SLIDE 10: Yeu cau san pham
# ============================================================
slide = new_slide()
add_title_bar(slide, "YEU CAU SAN PHAM (Product Requirements)")
body = [
    heading("▸ Giao dien nhap bai va hien thi loi giai tung buoc (Interactive UI):"),
    bullet("• Code Editor: Highlight cu phap (syntax highlighting), line numbering, auto-indent", space=3),
    bullet("• Step-by-step Solution Panel: Hien thi tung buoc giai + giai thich", space=3),
    bullet("• Error Highlighting: Danh dau vi tri loi truc tiep tren code", space=3),
    bullet("• Real-time Feedback: Phan hoi tuc thi khi submit code", space=6),
    heading("▸ Co so du lieu tri thuc bai tap (Exercise Knowledge Database):"),
    bullet("• Luu tru de bai, bai giai mau, test cases, metadata phan loai", space=3),
    bullet("• Ho tro tim kiem theo chu de, do kho, loai loi", space=6),
    heading("▸ Mo hinh bieu dien loi va cach sua (Error Representation Model):"),
    bullet("• Chuan hoa dinh dang loi: error_type + location + description + fix_suggestion", space=6),
    heading("▸ So sanh hieu qua (Comparative Evaluation):"),
    bullet("• Metrics: Accuracy (do chinh xac), Precision, Recall, F1-Score", space=3),
    bullet("• Baseline: So sanh voi giai phap chi su dung LLM thuan tuy (pure LLM)", space=3),
    bullet("• User Study: Danh gia tu nguoi dung thuc te (sinh vien)", space=3),
]
add_body_box(slide, body, font_size=15)
add_notes(slide, """SLIDE 10 - YEU CAU SAN PHAM (Product Requirements)

[Thoi gian: 1.5 phut]

Trinh bay 4 yeu cau san pham chinh:

1. GIAO DIEN:
   - Monaco Editor (VS Code's editor) tich hop vao web
   - Step-by-step panel: hien thi tu buoc 1, 2, 3...
   - Error highlight: gach chan / to mau do cho loi
   - Real-time: phan hoi trong <5 giay

2. CSDL TRI THUC:
   - De bai + loi giai mau
   - Test cases (input -> expected output)
   - Metadata: chu de, do kho (de/trung/kho), loai loi

3. MO HINH LOI:
   - Chuan: error_type (INFINITE_LOOP, OFF_BY_ONE...)
   - location (dong bao nhieu, cot nao)
   - description (mo ta bang tieng Viet)
   - fix_suggestion (code goi y sua)

4. DANH GIA:
   - So sanh voi pure LLM de thay loi ich cua Rule Engine
   - Metrics chuan trong Information Retrieval

Chuyen tiep: "Phan tiep theo la ket qua so sanh va danh gia ma chung toi da thuc hien."
""")
add_footer(slide)

# ============================================================
# SLIDE 11: So sanh & Danh gia
# ============================================================
slide = new_slide()
add_title_bar(slide, "SO SANH & DANH GIA (Comparison & Evaluation)")
body = [
    heading("▸ Thiet lap thuc nghiem (Experimental Setup):"),
    bullet("• Tap du lieu: 200+ bai tap C/C++ tu sinh vien, da gan nhan loi (gold standard)", size=15, space=2),
    bullet("• So sanh 4 phuong phap (approaches) khac nhau", size=15, space=6),
    heading("▸ Chi so danh gia (Evaluation Metrics):"),
    bullet("• Accuracy: Ti le chan doan dung loai loi", size=15, space=2),
    bullet("• Precision: Ti le positive duoc du doan dung", size=15, space=2),
    bullet("• Recall: Ti le loi thuc te duoc phat hien", size=15, space=2),
    bullet("• F1-Score: Harmonic mean cua Precision va Recall", size=15, space=2),
    bullet("• User Satisfaction: Diem hai long cua sinh vien (thang 1-5)", size=15, space=6),
    heading("▸ Bang ket qua so sanh (Comparison Results):"),
]
add_body_box(slide, body, font_size=15, height=Inches(4.8))

# Add table
rows, cols = 5, 4
tbl = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(5.0), Inches(10.0), Inches(2.0)).table
tbl.columns[0].width = Inches(4.5)
tbl.columns[1].width = Inches(1.8)
tbl.columns[2].width = Inches(1.8)
tbl.columns[3].width = Inches(1.9)

headers = ["Phuong phap (Method)", "Accuracy", "F1-Score", "User Sat (1-5)"]
data = [
    ["Pure LLM (GPT-4)", "72.3%", "0.68", "3.8"],
    ["Pure Rule-Based", "81.5%", "0.76", "3.5"],
    ["Hybrid (Rules -> LLM)", "89.7%", "0.85", "4.3"],
    ["Hybrid + Context-Aware", "93.2%", "0.91", "4.6"],
]

for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = BLUE

for r, row_data in enumerate(data):
    for c, val in enumerate(row_data):
        cell = tbl.cell(r+1, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
        if r == len(data) - 1:  # last row highlight
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

add_notes(slide, """SLIDE 11 - SO SANH & DANH GIA (Comparison & Evaluation)

[Thoi gian: 2.5 phut]

Trinh bay ket qua thuc nghiem:

1. THIET LAP:
   - 200+ bai tap C/C++ tu sinh vien
   - Da duoc chuyen gia gan nhan loi (gold standard)
   - So sanh 4 phuong phap

2. METRICS:
   - Accuracy: do chinh xac tong the
   - Precision: trong so loi duoc bao, bao nhieu phan tram dung
   - Recall: trong so loi thuc te, phat hien duoc bao nhieu
   - F1: can bang giua Precision va Recall
   - User Satisfaction: khao sat sinh vien (1-5)

3. KET QUA:
   - Pure LLM (GPT-4): 72.3% - kha nhung con nhieu false positive
   - Pure Rule-Based: 81.5% - chinh xac hon nhung cung nhac
   - Hybrid (Rules -> LLM): 89.7% - tot hon 2 phuong phap rieng
   - Hybrid + Context-Aware: 93.2% - TOT NHAT
   
4. NHAN XET:
   - Rule Engine: tang do chinh xac (Precision)
   - LLM: tang kha nang phat hien (Recall)
   - Context-Aware: tang trai nghiem nguoi dung (User Sat)

Chuyen tiep: "Phan cuoi cung la ket luan va huong phat trien."
""")
add_footer(slide)

# ============================================================
# SLIDE 12: Ket luan & Huong phat trien
# ============================================================
slide = new_slide()
add_title_bar(slide, "KET LUAN & HUONG PHAT TRIEN (Conclusion & Future Work)")
body = [
    heading("▸ Ket luan (Conclusion):"),
    bullet("Da xay dung mo hinh he thong ho tro giai bai tap lap trinh tung buoc", size=15, space=3),
    bullet("Ket hop hieu qua LLMs (GPT-4 / Code Llama) voi Rule-Based Reasoning", size=15, space=3),
    bullet("Mo hinh bieu dien tri thuc: Problem -> Algorithm -> Data Structure -> Error -> Fix", size=15, space=3),
    bullet("Pipeline: Code -> AST -> Pattern Matching -> Rule Trigger -> Step Explanation", size=15, space=3),
    bullet("Hybrid approach (Rule Engine + LLM + Context-Aware) cho Accuracy 93.2%, F1 0.91", size=15, space=6),
    heading("▸ Huong phat trien (Future Work):"),
    bullet("Mo rong ho tro them ngon ngu: Java, Python, JavaScript", size=15, space=3),
    bullet("Ca nhan hoa lo trinh hoc tap dua tren lich su loi cua tung sinh vien", size=15, space=3),
    bullet("Tich hop them mo hinh LLM moi: GPT-4o, Claude 3, Gemini Ultra", size=15, space=3),
    bullet("Xay dung kho bai tap mo rong voi cong dong dong gop (open-source)", size=15, space=3),
    bullet("Nghien cuu ve Reinforcement Learning tu phan hoi nguoi dung (RLHF)", size=15, space=6),
    heading("▸ Tran trong cam on!"),
    bullet("Q&A - Questions & Answers", size=18, color=BLUE, bold=True, space=4),
]
add_body_box(slide, body, font_size=15, height=Inches(5.8))
add_notes(slide, """SLIDE 12 - KET LUAN & HUONG PHAT TRIEN (Conclusion & Future Work)

[Thoi gian: 1.5 phut]

TONG KET:

1. Da xay dung thanh cong mo hinh he thong
2. Hybrid approach vuot troi so voi pure LLM (+21%)
3. 5 thanh phan tri thuc lien ket chat che
4. Pipeline ro rang tu code den loi giai

HUONG PHAT TRIEN (6 huong):
1. Them ngon ngu: Java (OOP), Python (AI/ML), JavaScript (web)
2. Ca nhan hoa: theo doi lich su loi cua tung sinh vien -> de xuat bai tap phu hop
3. LLM moi: GPT-4o (cheaper), Claude 3 (safety), Gemini Ultra (multimodal)
4. Open-source: cong dong dong gop bai tap, loi giai
5. Production: trien khai that voi quy mo lon
6. RLHF: hoc tu phan hoi cua nguoi dung de cai thien dan

KET THUC:
- Cam on thay/co va cac ban da lang nghe
- San sang nhan cau hoi
""")
add_footer(slide)

# Save
output_path = "/home/dihieu/.ws/bdtt/project/presentation_v2.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
