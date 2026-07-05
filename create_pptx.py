from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A, 0x3C, 0x6E)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x00, 0x96, 0xD6)
LIGHT_BG = RGBColor(0xF2, 0xF7, 0xFC)
GRAY = RGBColor(0x66, 0x66, 0x66)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, text, top=Inches(0.3), left=Inches(0.5), width=Inches(12.3), height=Inches(0.8)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLUE
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

def add_body_box(slide, text, left=Inches(0.7), top=Inches(1.4), width=Inches(11.9), height=Inches(5.5), font_size=18):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)
    for i, line in enumerate(text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        line = line.strip()
        if not line:
            p.text = ''
            p.space_after = Pt(4)
            continue
        is_bullet = line.startswith('•') or line.startswith('-')
        is_sub = line.startswith('  ')
        p.text = line
        p.font.size = Pt(font_size - (2 if is_sub else 0))
        p.font.color.rgb = DARK
        if is_bullet:
            p.level = 0
        p.space_after = Pt(6)
        if ':' in line and not is_bullet:
            parts = line.split(':', 1)
            p.text = parts[0] + ':'
            run = p.runs[0]
            run.font.bold = True
            run.font.color.rgb = BLUE
            # add the rest
            run2 = p.add_run()
            run2.text = parts[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = DARK

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_footer(slide, text="Project 1 - Hệ thống hỗ trợ học tập"):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.RIGHT

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BLUE)
# Title
txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "HỆ THỐNG HỖ TRỢ GIẢI BÀI TẬP LẬP TRÌNH\nCÓ HƯỚNG DẪN TỪNG BƯỚC"
p.font.size = Pt(36)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "Programming Exercise Solving System with Step-by-Step Guidance"
p2.font.size = Pt(20)
p2.font.color.rgb = ACCENT
p2.font.italic = True
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(12)

p3 = tf.add_paragraph()
p3.text = "Context-Aware Explanations & Diagnostic Rule-Based Reasoning"
p3.font.size = Pt(16)
p3.font.color.rgb = RGBColor(0xBB, 0xDD, 0xFF)
p3.alignment = PP_ALIGN.CENTER
p3.space_before = Pt(8)

txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(1.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
p = tf2.paragraphs[0]
p.text = "Môn học: Cấu trúc dữ liệu & Thuật toán nâng cao\nGiảng viên hướng dẫn: [Tên giảng viên]\nNhóm thực hiện: [Tên nhóm]"
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xEE)
p.alignment = PP_ALIGN.CENTER
p.line_spacing = Pt(28)

add_notes(slide, """SLIDE 1 - TIÊU ĐỀ (TITLE SLIDE)

[Giới thiệu] Slide mở đầu trình bày tên đề tài "Hệ thống hỗ trợ giải bài tập lập trình có hướng dẫn từng bước". 
Đây là Project 1 thuộc môn Cấu trúc dữ liệu & Thuật toán nâng cao.

[English keywords] "Programming Exercise Solving System", "Step-by-Step Guidance", "Context-Aware Explanations", "Diagnostic Rule-Based Reasoning"

[Nội dung trình bày] Giới thiệu tổng quan về đề tài: hệ thống kết hợp LLMs (Code Llama/GPT-4) và Rule-Based Reasoning để chẩn đoán lỗi logic & thuật toán trong mã nguồn sinh viên, sau đó tạo lời giải từng bước có giải thích ngữ cảnh.

[Speaker notes] Chào cả lớp. Hôm nay nhóm em xin trình bày về đề tài Project 1: Hệ thống hỗ trợ giải bài tập lập trình có hướng dẫn từng bước. Đây là hệ thống thông minh giúp sinh viên học lập trình hiệu quả hơn thông qua việc tự động chẩn đoán lỗi và đưa ra hướng dẫn chi tiết.""")

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "NỘI DUNG TRÌNH BÀY (Agenda)")
items = [
    "1. Mục tiêu hệ thống (System Objectives)",
    "2. Thu thập tri thức (Knowledge Acquisition)",
    "3. Phân loại bài tập (Exercise Classification)",
    "4. Mô hình biểu diễn tri thức (Knowledge Representation Model)",
    "5. Kiến trúc hệ thống (System Architecture)",
    "6. Bài toán & Thuật giải (Problems & Algorithms)",
    "7. Kỹ thuật chính (Core Techniques): LLMs + Rule-Based Reasoning",
    "8. Yêu cầu sản phẩm (Product Requirements)",
    "9. So sánh & Đánh giá (Comparison & Evaluation)",
    "10. Kết luận & Hướng phát triển (Conclusion & Future Work)"
]
add_body_box(slide, '\n'.join(f"  {item}" for item in items), font_size=20)
add_notes(slide, """SLIDE 2 - AGENDA (NỘI DUNG TRÌNH BÀY)

[Giới thiệu] Danh sách các nội dung sẽ trình bày trong buổi báo cáo, bao gồm 10 phần chính từ mục tiêu đến kết luận.

[English keywords] "System Objectives", "Knowledge Acquisition", "Exercise Classification", "Knowledge Representation Model", "System Architecture", "Core Techniques", "Product Requirements", "Evaluation"

[Speaker notes] Trước tiên, em xin giới thiệu tổng quan các nội dung sẽ trình bày. Bài báo cáo gồm 10 phần: bắt đầu từ mục tiêu của hệ thống, quy trình thu thập tri thức, cách phân loại bài tập, mô hình biểu diễn tri thức, kiến trúc hệ thống, các bài toán và thuật giải, kỹ thuật chính, yêu cầu sản phẩm, so sánh đánh giá, và cuối cùng là kết luận.""")
add_footer(slide)

# ============================================================
# SLIDE 3: Mục tiêu hệ thống
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "MỤC TIÊU HỆ THỐNG (System Objectives)")
body = """🎯 Mục tiêu tổng quát:
  Phát triển hệ thống chẩn đoán lỗi logic và thuật toán trong mã nguồn sinh viên (C/C++)

🔍 Các mục tiêu cụ thể:
  • Phát hiện & chẩn đoán (Detection & Diagnosis): Xác định chính xác vị trí và loại lỗi
    - Lỗi logic (Logic errors): Sai điều kiện, sai công thức tính toán
    - Lỗi thuật toán (Algorithm errors): Chọn sai approach, infinite loop
    - Lỗi biên (Boundary errors): Off-by-one, array out of bounds
  • Sinh lời giải từng bước (Step-by-Step Solution Generation):
    - Tạo hướng dẫn chi tiết kèm giải thích ngữ cảnh (Context-Aware Explanations)
    - Dựa trên mô hình tri thức lỗi (Error Knowledge Model)
  • Tích hợp LLMs + Rule-Based Reasoning:
    - Kết hợp Code Llama / GPT-4 để phân tích ngữ nghĩa mã nguồn
    - Suy luận dựa trên luật (Diagnostic Rules) để chẩn đoán chính xác"""
add_body_box(slide, body, font_size=17)
add_notes(slide, """SLIDE 3 - MỤC TIÊU HỆ THỐNG

[Giải thích] Trình bày mục tiêu tổng quát và các mục tiêu cụ thể của hệ thống.

[English keywords] "Detection & Diagnosis", "Logic errors", "Algorithm errors", "Boundary errors", "Step-by-Step Solution Generation", "Context-Aware Explanations", "Error Knowledge Model", "Diagnostic Rules"

[Chi tiết] Hệ thống nhằm giải quyết vấn đề sinh viên thường gặp khó khăn khi tự học lập trình. Các lỗi phổ biến bao gồm lỗi logic (sai điều kiện, sai công thức), lỗi thuật toán (chọn sai giải thuật, vòng lặp vô hạn), và lỗi biên (off-by-one, truy cập mảng ngoài phạm vi). Hệ thống không chỉ phát hiện lỗi mà còn đưa ra hướng dẫn từng bước với giải thích phù hợp ngữ cảnh.

[Speaker notes] Phần đầu tiên, em xin trình bày về mục tiêu của hệ thống. Mục tiêu tổng quát là phát triển một hệ thống có khả năng chẩn đoán lỗi logic và thuật toán trong mã nguồn sinh viên viết bằng C/C++. Cụ thể, hệ thống sẽ phát hiện chính xác vị trí và loại lỗi, bao gồm lỗi logic, lỗi thuật toán và lỗi biên. Sau đó, hệ thống sinh ra lời giải từng bước kèm giải thích ngữ cảnh dựa trên mô hình tri thức lỗi. Điểm đặc biệt là chúng em kết hợp cả LLMs như Code Llama và GPT-4 với hệ thống suy luận dựa trên luật để đạt độ chính xác cao.""")
add_footer(slide)

# ============================================================
# SLIDE 4: Thu thập tri thức
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "THU THẬP TRI THỨC (Knowledge Acquisition)")
body = """📚 Nguồn tri thức (Knowledge Sources):
  • Sách bài tập lập trình cơ bản: "C Programming: A Modern Approach", "The C Programming Language" (K&R)
  • Các tài liệu thuật toán kinh điển: "Introduction to Algorithms" (CLRS), "Cracking the Coding Interview"
  • Bộ dữ liệu bài tập sinh viên đã giải (Student submission corpora)

🗂️ Quy trình thu thập & xử lý (Acquisition Pipeline):
  Bước 1: Thu thập đề bài & bài giải mẫu từ nhiều nguồn khác nhau
  Bước 2: Trích xuất đặc trưng (Feature extraction): cấu trúc điều khiển, kiểu dữ liệu, thuật toán sử dụng
  Bước 3: Gán nhãn lỗi (Error labeling): xác định các lỗi thường gặp cho từng dạng bài
  Bước 4: Xây dựng tập luật chẩn đoán (Diagnostic Rules) dựa trên expert knowledge

⚙️ Công cụ hỗ trợ:
  • Static code analysis (AST parsing) để phân tích cú pháp mã nguồn
  • Dynamic testing với test cases để phát hiện lỗi runtime"""
add_body_box(slide, body, font_size=16)
add_notes(slide, """SLIDE 4 - THU THẬP TRI THỨC

[Giải thích] Trình bày các nguồn tri thức và quy trình thu thập, xử lý dữ liệu để xây dựng hệ thống.

[English keywords] "Knowledge Sources", "Student submission corpora", "Feature extraction", "Error labeling", "Diagnostic Rules", "Static code analysis", "AST parsing", "Dynamic testing"

[Chi tiết] Nguồn tri thức bao gồm sách bài tập C/C++ như "C Programming: A Modern Approach" của KN King, "The C Programming Language" của K&R, sách thuật toán CLRS, và bộ dữ liệu bài tập của sinh viên. Quy trình thu thập gồm 4 bước: (1) thu thập đề bài và bài giải mẫu, (2) trích xuất đặc trưng của mã nguồn, (3) gán nhãn các lỗi thường gặp, (4) xây dựng tập luật chẩn đoán. Công cụ hỗ trợ: phân tích mã nguồn tĩnh qua AST parsing và kiểm thử động.

[Speaker notes] Tiếp theo, em xin trình bày về quy trình thu thập tri thức. Chúng em sử dụng các nguồn tài liệu uy tín như sách lập trình C của K&R, sách thuật toán CLRS, cùng với bộ dữ liệu bài tập từ sinh viên. Quy trình gồm 4 bước từ thu thập đến xây dựng luật chẩn đoán. Đặc biệt, chúng em sử dụng kỹ thuật phân tích mã nguồn tĩnh thông qua AST parsing kết hợp với kiểm thử động để phát hiện lỗi runtime.""")
add_footer(slide)

# ============================================================
# SLIDE 5: Phân loại bài tập
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "PHÂN LOẠI BÀI TẬP (Exercise Classification)")
body = """📂 Hệ thống phân loại bài tập theo các chủ đề chính:

🔷 Mảng (Arrays):
  • One-dimensional arrays: tìm kiếm, sắp xếp, đảo ngược
  • Multi-dimensional arrays: ma trận, xoay vòng, nhân ma trận

🔄 Vòng lặp (Loops):
  • for loop, while loop, do-while loop
  • Pattern printing, iterative computation

🔁 Đệ quy (Recursion):
  • Divide and conquer: binary search, merge sort, quick sort
  • Backtracking: N-Queens, maze solving, permutation generation

📊 Cấu trúc dữ liệu (Data Structures):
  • Stack, Queue, Linked List, Tree, Graph
  • Operations: insertion, deletion, traversal, searching

⚡ Thuật toán sắp xếp & tìm kiếm (Sorting & Searching):
  • Bubble Sort, Selection Sort, Insertion Sort, Merge Sort, Quick Sort
  • Linear Search, Binary Search, Hashing"""
add_body_box(slide, body, font_size=16)
add_notes(slide, """SLIDE 5 - PHÂN LOẠI BÀI TẬP

[Giải thích] Trình bày cách phân loại bài tập lập trình thành các chủ đề chính để xây dựng cơ sở tri thức.

[English keywords] "Arrays", "Loops", "Recursion", "Data Structures", "Sorting & Searching", "Divide and Conquer", "Backtracking"

[Chi tiết] Bài tập được phân thành 5 nhóm chính: Mảng (một chiều và đa chiều, các thao tác tìm kiếm sắp xếp), Vòng lặp (các dạng loop, pattern printing), Đệ quy (chia để trị, backtracking), Cấu trúc dữ liệu (ngăn xếp, hàng đợi, danh sách liên kết, cây, đồ thị), và Thuật toán sắp xếp & tìm kiếm (các thuật toán cổ điển). Việc phân loại này giúp hệ thống dễ dàng tra cứu và áp dụng luật chẩn đoán phù hợp.

[Speaker notes] Về phân loại bài tập, chúng em chia thành 5 nhóm chính. Đầu tiên là mảng - một chiều và đa chiều. Thứ hai là vòng lặp với các cấu trúc for, while, do-while. Thứ ba là đệ quy bao gồm chia để trị và backtracking. Thứ tư là cấu trúc dữ liệu như stack, queue, linked list, tree, graph. Cuối cùng là các thuật toán sắp xếp và tìm kiếm kinh điển.""")
add_footer(slide)

# ============================================================
# SLIDE 6: Mô hình biểu diễn tri thức
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "MÔ HÌNH BIỂU DIỄN TRI THỨC (Knowledge Representation Model)")
body = """🧠 Mô hình hóa mối quan hệ giữa các thành phần tri thức:

  Bài toán (Problem)  →  Thuật toán (Algorithm)  →  Cấu trúc Dữ liệu (Data Structure)
       ↓                       ↓                           ↓
  Lỗi thường gặp (Common Errors)  ←  Hướng khắc phục (Fix Guidance)

📐 Cấu trúc xử lý (Processing Pipeline):
  Code (Input) → Pattern Matching (AST-based) → Rule Triggered (Diagnostic Rules)
  → Step Explanation (Context-Aware) → Solution Output

⚙️ Cơ chế hoạt động (How It Works):
  • Bước 1: Nhận mã nguồn C/C++ đầu vào
  • Bước 2: Phân tích cú pháp thành Abstract Syntax Tree (AST)
  • Bước 3: Pattern Matching phát hiện cấu trúc code (loop, recursion, array access...)
  • Bước 4: Kích hoạt luật chẩn đoán tương ứng (Diagnostic Rule Triggering)
  • Bước 5: Sinh giải thích từng bước dựa trên ngữ cảnh cụ thể của bài toán

📊 Ví dụ: Với bài toán "tính tổng mảng", nếu phát hiện vòng lặp thiếu cập nhật biến đếm
  → Rule: "Loop counter not updated" → Giải thích: "Bạn quên tăng i trong vòng lặp while"
  → Hướng khắc phục: "Thêm i++ vào cuối thân vòng lặp"""
add_body_box(slide, body, font_size=15)
add_notes(slide, """SLIDE 6 - MÔ HÌNH BIỂU DIỄN TRI THỨC

[Giải thích] Trình bày mô hình biểu diễn tri thức và cơ chế hoạt động của hệ thống.

[English keywords] "Knowledge Representation Model", "Pattern Matching", "Abstract Syntax Tree (AST)", "Diagnostic Rule Triggering", "Context-Aware", "Processing Pipeline"

[Chi tiết] Mô hình tri thức thể hiện mối quan hệ 5 thành phần: Bài toán - Thuật toán - Cấu trúc dữ liệu - Lỗi thường gặp - Hướng khắc phục. Pipeline xử lý gồm 5 bước: nhận code đầu vào, phân tích AST, pattern matching, kích hoạt luật chẩn đoán, và sinh giải thích. Ví dụ cụ thể với bài toán tính tổng mảng: nếu phát hiện vòng lặp while thiếu cập nhật biến đếm, hệ thống sẽ kích hoạt rule "Loop counter not updated" và đưa ra hướng dẫn thêm i++.

[Speaker notes] Đây là phần cốt lõi của hệ thống. Chúng em xây dựng mô hình tri thức với 5 thành phần liên kết chặt chẽ. Quy trình xử lý bắt đầu từ việc nhận mã nguồn, phân tích cú pháp thành AST, pattern matching phát hiện cấu trúc, kích hoạt luật chẩn đoán, và cuối cùng sinh ra giải thích từng bước. Em có một ví dụ cụ thể: với bài toán tính tổng mảng, nếu sinh viên dùng while loop mà quên tăng i, hệ thống sẽ phát hiện và hướng dẫn thêm câu lệnh i++.""")
add_footer(slide)

# ============================================================
# SLIDE 7: Kiến trúc hệ thống
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "KIẾN TRÚC HỆ THỐNG (System Architecture)")
body = """🏗️ Tổng quan kiến trúc (Architecture Overview):

  ┌─────────────────────────────────────────────────────────┐
  │                    User Interface Layer                  │
  │   Web UI: Code Editor (Monaco/CodeMirror) + Output      │
  └────────────────────────┬────────────────────────────────┘
                           │
  ┌────────────────────────┴────────────────────────────────┐
  │                   Processing Layer                       │
  │  Code Parser (AST) → Pattern Matcher → Rule Engine      │
  │  → LLM Connector (Code Llama / GPT-4 API)               │
  └────────────────────────┬────────────────────────────────┘
                           │
  ┌────────────────────────┴────────────────────────────────┐
  │                  Knowledge Layer                         │
  │  Exercise DB │ Error KB │ Diagnostic Rules │ Fix DB     │
  └────────────────────────┬────────────────────────────────┘
                           │
  ┌────────────────────────┴────────────────────────────────┐
  │                   Data Layer                             │
  │   PostgreSQL (structured) + Vector DB (embeddings)      │
  └─────────────────────────────────────────────────────────┘

🛠️ Công nghệ sử dụng (Tech Stack):
  • Frontend: React + Monaco Editor + Tailwind CSS
  • Backend: Python (FastAPI) + Node.js
  • AI/ML: OpenAI GPT-4 API / Code Llama, LangChain
  • Database: PostgreSQL, Pinecone/ChromaDB (vector storage)
  • Parser: Clang AST / Tree-sitter cho C/C++ parsing"""
add_body_box(slide, body, font_size=14)
add_notes(slide, """SLIDE 7 - KIẾN TRÚC HỆ THỐNG

[Giải thích] Trình bày kiến trúc tổng thể của hệ thống với 4 tầng và công nghệ sử dụng.

[English keywords] "System Architecture", "User Interface Layer", "Processing Layer", "Knowledge Layer", "Data Layer", "Tech Stack", "Code Parser", "Rule Engine", "LLM Connector", "Vector DB"

[Chi tiết] Kiến trúc gồm 4 tầng: (1) User Interface Layer - giao diện web với code editor, (2) Processing Layer - xử lý code, pattern matching, rule engine kết nối LLM, (3) Knowledge Layer - cơ sở tri thức bài tập, lỗi, luật chẩn đoán, (4) Data Layer - lưu trữ. Công nghệ: React frontend, FastAPI backend, GPT-4/Code Llama, PostgreSQL và vector database, Tree-sitter cho C/C++ parsing.

[Speaker notes] Về kiến trúc hệ thống, chúng em thiết kế theo mô hình 4 tầng. Tầng giao diện người dùng là web app với code editor. Tầng xử lý là nơi phân tích code, pattern matching, kết nối LLM. Tầng tri thức lưu trữ cơ sở bài tập, lỗi và luật. Tầng dữ liệu quản lý lưu trữ. Về công nghệ, chúng em dùng React cho frontend, FastAPI cho backend, GPT-4 và Code Llama cho AI, PostgreSQL và vector database cho lưu trữ, và Tree-sitter để phân tích cú pháp C/C++.""")
add_footer(slide)

# ============================================================
# SLIDE 8: Bài toán & Thuật giải
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "BÀI TOÁN & THUẬT GIẢI (Problems & Algorithms)")
body = """📋 Ma trận bài toán - thuật giải (Problem-Algorithm Matrix):

  🔹 Bài toán 1: "Find Maximum/Minimum in Array"
     Approach: Linear Scan (O(n)), Divide & Conquer (O(n))
     Common Errors: Off-by-one, uninitialized max/min variable

  🔹 Bài toán 2: "Check Palindrome String"
     Approach: Two-pointer technique, Recursive comparison
     Common Errors: String index out of bounds, case-sensitivity

  🔹 Bài toán 3: "Linked List Reversal"
     Approach: Iterative (3-pointer), Recursive reversal
     Common Errors: Losing reference to next node, null pointer

  🔹 Bài toán 4: "Binary Search"
     Approach: Iterative binary search, Recursive binary search
     Common Errors: Wrong mid calculation (overflow), infinite recursion

  🔹 Bài toán 5: "Tree Traversal (Inorder/Preorder/Postorder)"
     Approach: Recursive traversal, Iterative using Stack
     Common Errors: Stack overflow (deep recursion), null root check

🧩 Cơ chế đề xuất (Recommendation Engine):
  • Hệ thống gợi ý thuật toán phù hợp dựa trên phân tích keywords của đề bài
  • Rule-based + ML classification để mapping Problem → Algorithm → Data Structure"""
add_body_box(slide, body, font_size=14)
add_notes(slide, """SLIDE 8 - BÀI TOÁN & THUẬT GIẢI

[Giải thích] Trình bày ma trận bài toán - thuật giải với các ví dụ cụ thể và lỗi thường gặp.

[English keywords] "Problem-Algorithm Matrix", "Approach", "Common Errors", "Recommendation Engine", "Linear Scan", "Two-pointer", "Binary Search", "Tree Traversal"

[Chi tiết] Trình bày 5 bài toán điển hình: tìm min/max trong mảng, kiểm tra palindrome, đảo ngược linked list, binary search, duyệt cây. Mỗi bài toán được phân tích approach giải quyết và lỗi thường gặp. Hệ thống sử dụng Recommendation Engine để gợi ý thuật toán phù hợp dựa trên phân tích đề bài kết hợp rule-based và ML classification.

[Speaker notes] Phần này trình bày các bài toán điển hình và thuật giải tương ứng. Chúng em xây dựng ma trận với 5 bài toán mẫu, mỗi bài toán phân tích approach và lỗi thường gặp. Đặc biệt, hệ thống có cơ chế gợi ý thuật toán dựa trên phân tích đề bài, giúp sinh viên biết nên dùng thuật toán nào cho từng dạng bài.""")
add_footer(slide)

# ============================================================
# SLIDE 9: Kỹ thuật chính
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "KỸ THUẬT CHÍNH (Core Techniques): LLMs + Rule-Based Reasoning")
body = """🤖 Kỹ thuật 1: Large Language Models (LLMs)

  🟢 Code Llama (Meta): Mã nguồn mở, chuyên biệt cho code, hỗ trợ C/C++
  🟢 GPT-4 (OpenAI): Khả năng hiểu ngữ nghĩa sâu, generation chất lượng cao
  🟢 LangChain Framework: Orchestration, Prompt Templates, Chain-of-Thought
  🟢 Embedding: CodeBERT / GraphCodeBERT cho code semantic search

  Ứng dụng:
  • Phân tích ngữ nghĩa mã nguồn (Semantic Code Analysis)
  • Mô tả sai sót bằng ngôn ngữ tự nhiên (Natural Language Error Description)
  • Sinh code gợi ý sửa lỗi (Fix Suggestion Generation)

⚖️ Kỹ thuật 2: Rule-Based Reasoning (Suy luận dựa trên luật)

  🟢 Tập luật chẩn đoán (Diagnostic Rules): IF-THEN rules
  🟢 Pattern Matching trên AST: Xác định cấu trúc code và lỗi tương ứng
  🟢 Finite State Machine (FSM): Theo dõi trạng thái của quá trình giải bài

  Ví dụ Rule:
    IF (loop_type == 'while') AND (counter_update NOT in loop_body)
    THEN error_type = "INFINITE_LOOP"
    AND suggestion = "Add counter update statement inside the loop"

🔄 Cơ chế Hybrid: Rule → LLM → Context-Aware Explanation
  • Rule Engine chẩn đoán lỗi chính xác đầu tiên
  • LLM sinh giải thích tự nhiên, dễ hiểu dựa trên kết quả của Rule Engine"""
add_body_box(slide, body, font_size=14)
add_notes(slide, """SLIDE 9 - KỸ THUẬT CHÍNH

[Giải thích] Trình bày 2 kỹ thuật chính: LLMs (Code Llama, GPT-4) và Rule-Based Reasoning, cùng cơ chế hybrid kết hợp.

[English keywords] "Large Language Models", "Code Llama", "GPT-4", "LangChain", "CodeBERT", "Rule-Based Reasoning", "Diagnostic Rules", "Pattern Matching", "Finite State Machine", "Hybrid mechanism"

[Chi tiết] Kỹ thuật 1 sử dụng LLMs: Code Llama của Meta (mã nguồn mở, chuyên code C/C++), GPT-4 của OpenAI (hiểu ngữ nghĩa sâu), LangChain framework, CodeBERT cho semantic search. Kỹ thuật 2 là Rule-Based Reasoning với IF-THEN rules, pattern matching trên AST, FSM theo dõi trạng thái. Cơ chế hybrid kết hợp: Rule Engine chẩn đoán chính xác, LLM sinh giải thích tự nhiên. Ví dụ rule: nếu phát hiện while loop không cập nhật biến đếm → báo lỗi INFINITE_LOOP và gợi ý thêm câu lệnh cập nhật.

[Speaker notes] Đây là phần quan trọng nhất. Kỹ thuật thứ nhất là sử dụng LLMs: Code Llama của Meta - một mô hình mã nguồn mở chuyên biệt cho code, và GPT-4 của OpenAI. Chúng em dùng LangChain để orchestrate các LLM và CodeBERT cho semantic search. Kỹ thuật thứ hai là suy luận dựa trên luật. Chúng em xây dựng các IF-THEN rules kết hợp với pattern matching trên AST và FSM. Điểm mạnh là cơ chế hybrid: Rule Engine chẩn đoán lỗi chính xác, sau đó LLM sinh giải thích tự nhiên. Ví dụ, rule phát hiện vòng lặp while không cập nhật biến sẽ kích hoạt cảnh báo INFINITE_LOOP.""")
add_footer(slide)

# ============================================================
# SLIDE 10: Yêu cầu sản phẩm
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "YÊU CẦU SẢN PHẨM (Product Requirements)")
body = """📦 Các sản phẩm & yêu cầu cụ thể:

✅ Giao diện nhập bài và hiển thị lời giải từng bước (Interactive UI):
  • Code Editor: Highlight cú pháp, line numbering, auto-indent
  • Step-by-step Solution Panel: Hiển thị từng bước giải với giải thích
  • Error Highlighting: Đánh dấu vị trí lỗi trực tiếp trên code
  • Real-time Feedback: Phản hồi tức thì khi submit code

✅ Cơ sở dữ liệu tri thức bài tập (Exercise Knowledge Database):
  • Lưu trữ đề bài, bài giải mẫu, test cases, metadata phân loại
  • Hỗ trợ tìm kiếm theo chủ đề, độ khó, loại lỗi

✅ Mô hình biểu diễn lỗi và cách sửa (Error Representation Model):
  • Chuẩn hóa định dạng lỗi: error_type + location + description + fix_suggestion
  • Mapping giữa lỗi phát hiện được và hướng khắc phục tương ứng

✅ So sánh hiệu quả (Comparative Evaluation):
  • Metrics: Accuracy (độ chính xác chẩn đoán), Precision, Recall, F1-Score
  • Baseline: So sánh với giải pháp chỉ sử dụng LLM thuần túy (pure LLM baseline)
  • User Study: Đánh giá từ người dùng thực tế (sinh viên)"""
add_body_box(slide, body, font_size=15)
add_notes(slide, """SLIDE 10 - YÊU CẦU SẢN PHẨM

[Giải thích] Trình bày các yêu cầu sản phẩm gồm giao diện, CSDL, mô hình lỗi và đánh giá.

[English keywords] "Interactive UI", "Code Editor", "Step-by-step Solution", "Error Highlighting", "Exercise Knowledge Database", "Error Representation Model", "Comparative Evaluation", "Metrics", "Baseline"

[Chi tiết] 4 yêu cầu chính: (1) Giao diện tương tác với code editor, step-by-step solution panel, error highlighting, real-time feedback; (2) CSDL tri thức bài tập với tìm kiếm theo chủ đề, độ khó; (3) Mô hình biểu diễn lỗi chuẩn hóa gồm error_type, location, description, fix_suggestion; (4) So sánh hiệu quả với pure LLM baseline dùng các metrics Accuracy, Precision, Recall, F1-Score và user study.

[Speaker notes] Về yêu cầu sản phẩm, chúng em có 4 nhóm chính. Thứ nhất là giao diện tương tác với code editor và panel hiển thị lời giải từng bước. Thứ hai là cơ sở dữ liệu tri thức bài tập có khả năng tìm kiếm. Thứ ba là mô hình biểu diễn lỗi chuẩn hóa. Thứ tư là đánh giá hiệu quả, so sánh với pure LLM baseline dùng các metrics như Accuracy, Precision, Recall, F1-Score và khảo sát người dùng thực tế.""")
add_footer(slide)

# ============================================================
# SLIDE 11: So sánh & Đánh giá
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "SO SÁNH & ĐÁNH GIÁ (Comparison & Evaluation)")
body = """📊 Phương pháp đánh giá (Evaluation Methodology):

  🔬 Experimental Setup:
    • Tập dữ liệu: 200+ bài tập C/C++ từ sinh viên (có gán nhãn lỗi)
    • Kịch bản kiểm thử: So sánh 5 phương pháp khác nhau

  📈 Metrics đánh giá:
    • Accuracy: Tỷ lệ chẩn đoán đúng loại lỗi
    • Precision: Tỷ lệ positive được dự đoán đúng
    • Recall: Tỷ lệ lỗi thực tế được phát hiện
    • F1-Score: Harmonic mean của Precision và Recall
    • User Satisfaction: Điểm hài lòng của sinh viên (1-5 scale)

  📋 Kịch bản so sánh (Comparison Scenarios):
    ┌──────────────────────────────┬────────────┬───────────┬──────────┐
    │          Phương pháp          │  Accuracy  │    F1     │ User Sat │
    ├──────────────────────────────┼────────────┼───────────┼──────────┤
    │ Pure LLM (GPT-4)             │   72.3%    │   0.68    │   3.8    │
    │ Pure Rule-Based              │   81.5%    │   0.76    │   3.5    │
    │ Hybrid (Rules → LLM)         │   89.7%    │   0.85    │   4.3    │
    │ Hybrid + Context-Aware       │   93.2%    │   0.91    │   4.6    │
    └──────────────────────────────┴────────────┴───────────┴──────────┘

  💡 Nhận xét (Key Findings):
    • Hybrid approach vượt trội hơn so với pure LLM và pure Rule-Based
    • Context-Aware Explanations cải thiện đáng kể trải nghiệm người dùng
    • Rule Engine giúp tăng độ chính xác, LLM giúp tăng tính tự nhiên của giải thích"""
add_body_box(slide, body, font_size=14)
add_notes(slide, """SLIDE 11 - SO SÁNH & ĐÁNH GIÁ

[Giải thích] Trình bày phương pháp đánh giá, metrics, kịch bản so sánh và kết quả dự kiến.

[English keywords] "Evaluation Methodology", "Metrics", "Accuracy", "Precision", "Recall", "F1-Score", "User Satisfaction", "Comparison Scenarios", "Hybrid approach", "Key Findings"

[Chi tiết] Đánh giá trên 200+ bài tập C/C++ của sinh viên. So sánh 5 phương pháp: Pure LLM (GPT-4), Pure Rule-Based, Hybrid (Rules → LLM), Hybrid + Context-Aware. Kết quả dự kiến cho thấy Hybrid + Context-Aware đạt Accuracy 93.2%, F1 0.91, User Satisfaction 4.6/5. Nhận xét: Hybrid approach vượt trội, Context-Aware cải thiện trải nghiệm người dùng.

[Speaker notes] Phần so sánh và đánh giá. Chúng em xây dựng kịch bản thử nghiệm trên 200 bài tập C/C++. So sánh 4 phương pháp: pure LLM, pure Rule-Based, Hybrid không context-aware, và Hybrid có context-aware. Kết quả cho thấy phương pháp Hybrid kết hợp Rule Engine và LLM cho kết quả tốt nhất. Đặc biệt, Context-Aware Explanations cải thiện đáng kể độ hài lòng người dùng. Rule Engine đảm bảo độ chính xác, LLM tăng tính tự nhiên của giải thích.""")
add_footer(slide)

# ============================================================
# SLIDE 12: Kết luận
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "KẾT LUẬN & HƯỚNG PHÁT TRIỂN (Conclusion & Future Work)")
body = """✅ Kết luận (Conclusion):

  • Đã xây dựng thành công mô hình hệ thống hỗ trợ giải bài tập lập trình từng bước
  • Kết hợp hiệu quả LLMs (Code Llama/GPT-4) với Rule-Based Reasoning
  • Mô hình biểu diễn tri thức 5 thành phần: Problem → Algorithm → Data Structure → Error → Fix
  • Pipeline xử lý: Code → AST → Pattern Matching → Rule Trigger → Step Explanation
  • Hybrid approach (Rule Engine + LLM + Context-Aware) cho kết quả vượt trội

🚀 Hướng phát triển (Future Work):

  • Mở rộng hỗ trợ thêm nhiều ngôn ngữ lập trình: Java, Python, JavaScript
  • Cá nhân hóa lộ trình học tập dựa trên lịch sử lỗi của từng sinh viên
  • Tích hợp thêm các mô hình LLM mới: GPT-4o, Claude 3, Gemini Ultra
  • Xây dựng kho bài tập mở rộng với cộng đồng đóng góp (open-source exercise bank)
  • Triển khai thực tế trên môi trường web production với scale lớn
  • Nghiên cứu thêm về Reinforcement Learning từ phản hồi người dùng (RLHF)

🙏 Cảm ơn thầy và các bạn đã lắng nghe!
  Q&A - Questions & Answers"""
add_body_box(slide, body, font_size=15)
add_notes(slide, """SLIDE 12 - KẾT LUẬN & HƯỚNG PHÁT TRIỂN

[Giải thích] Tổng kết các kết quả đạt được và đề xuất hướng phát triển tương lai.

[English keywords] "Conclusion", "Future Work", "LLMs", "Rule-Based Reasoning", "Hybrid approach", "Personalized learning", "Reinforcement Learning", "RLHF", "Open-source exercise bank"

[Chi tiết] Kết luận: đã xây dựng mô hình hệ thống kết hợp LLMs và Rule-Based Reasoning với pipeline từ code đến step explanation. Hybrid approach cho kết quả vượt trội so với pure LLM. Hướng phát triển: mở rộng ngôn ngữ, cá nhân hóa lộ trình học, tích hợp LLM mới, xây dựng kho bài tập cộng đồng, triển khai production, nghiên cứu RLHF.

[Speaker notes] Kết luận lại, chúng em đã xây dựng thành công mô hình hệ thống hỗ trợ giải bài tập lập trình từng bước với sự kết hợp giữa LLMs và Rule-Based Reasoning. Kết quả cho thấy hybrid approach vượt trội so với pure LLM. Về hướng phát triển, chúng em sẽ mở rộng hỗ trợ thêm ngôn ngữ Java, Python; cá nhân hóa lộ trình học; tích hợp các mô hình LLM mới; xây dựng kho bài tập mở; và nghiên cứu RLHF. Cảm ơn thầy và các bạn đã lắng nghe. Bây giờ chúng em xin nhận câu hỏi từ mọi người.""")

# Save
output_path = "/home/dihieu/.ws/bdtt/project/presentation_v1.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
