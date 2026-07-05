from pptx import Presentation
import json

prs = Presentation("/home/dihieu/.ws/bdtt/project/presentation_v1.pptx")
slides_data = []

for i, slide in enumerate(prs.slides, 1):
    slide_info = {"slide_number": i, "texts": [], "notes": ""}
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    slide_info["texts"].append(t)
    try:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        slide_info["notes"] = notes
    except:
        pass
    slides_data.append(slide_info)

print(json.dumps(slides_data, ensure_ascii=False, indent=2))
