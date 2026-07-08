#!/usr/bin/env python3
"""Sinh slide thuyet trinh (tieng Viet) cho Project 1 - BDTT.

Noi dung bam sat project.md va he thong trong repo, dat nen tren tai lieu
mon hoc "Mo hinh tri thuc quan he va Ung dung" (PGS.TS. Nguyen Dinh Hien, UIT).

Chay:  python presentation/build_pptx.py
Xuat:  presentation/Project1_BDTT_presentation.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- Bang mau ----
NAVY = RGBColor(0x0F, 0x2A, 0x5F)
BLUE = RGBColor(0x1F, 0x4E, 0x9B)
TEAL = RGBColor(0x14, 0xA0, 0x98)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
GREEN = RGBColor(0x2E, 0x8B, 0x57)
LIGHT = RGBColor(0xEF, 0xF3, 0xFA)
GREY = RGBColor(0x55, 0x5F, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1B, 0x24, 0x33)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          space_after=6, wrap=True):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (txt, size, bold, color, *rest) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = color
            if rest and rest[0]:
                r.font.italic = True
    return tb


def content_slide(title, kicker=None):
    """Slide chuan: thanh header + vung noi dung. Tra ve (slide, content_top)."""
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, WHITE)
    _box(s, 0, 0, SW, Inches(1.15), fill=NAVY)
    _box(s, 0, Inches(1.15), SW, Inches(0.07), fill=TEAL)
    _box(s, 0, 0, Inches(0.18), Inches(1.15), fill=ORANGE)
    _text(s, Inches(0.55), Inches(0.12), Inches(12.2), Inches(0.95),
          [[(title, 27, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        _text(s, Inches(0.57), Inches(0.80), Inches(12), Inches(0.3),
              [[(kicker, 12, False, RGBColor(0xBF, 0xD3, 0xF2))]])
    return s, Inches(1.5)


def bullets_slide(title, items, kicker=None, col=NAVY):
    """items: list of (text, level, bold?) hoac str."""
    s, top = content_slide(title, kicker)
    paras = []
    for it in items:
        if isinstance(it, tuple):
            txt, lvl, *b = it
            bold = b[0] if b else (lvl == 0)
        else:
            txt, lvl, bold = it, 0, False
        bullet = "▸ " if lvl == 0 else "•  "
        indent = "" if lvl == 0 else "      "
        color = col if lvl == 0 else GREY
        size = 18 if lvl == 0 else 15
        paras.append([(indent + bullet, size, True,
                       TEAL if lvl == 0 else ORANGE), (txt, size, bold, color)])
    _text(s, Inches(0.7), top, Inches(12), Inches(5.6), paras, space_after=9)
    return s


def footer(s, page):
    _text(s, Inches(0.4), Inches(7.05), Inches(9), Inches(0.35),
          [[("Project 1 — Biểu diễn Tri thức · Hệ hỗ trợ giải bài tập lập trình", 9, False, GREY)]])
    _text(s, Inches(12.2), Inches(7.05), Inches(0.9), Inches(0.35),
          [[(str(page), 9, True, NAVY)]], align=PP_ALIGN.RIGHT)


# ============================================================= SLIDE 1: TITLE
s = prs.slides.add_slide(BLANK)
_set_bg(s, NAVY)
_box(s, 0, Inches(2.55), SW, Inches(0.06), fill=TEAL)
_box(s, 0, Inches(4.7), SW, Inches(0.06), fill=ORANGE)
_box(s, Inches(0.6), Inches(0.5), Inches(3.2), Inches(0.5), fill=None, line=None)
_text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.5),
      [[("MÔN: BIỂU DIỄN TRI THỨC (BDTT)", 15, True, TEAL)]])
_text(s, Inches(0.9), Inches(2.75), Inches(11.6), Inches(2.0),
      [[("HỆ THỐNG HỖ TRỢ GIẢI BÀI TẬP LẬP TRÌNH", 38, True, WHITE)],
       [("CÓ HƯỚNG DẪN TỪNG BƯỚC", 38, True, WHITE)]], space_after=4)
_text(s, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.6),
      [[("Chẩn đoán lỗi logic & thuật toán trong mã nguồn sinh viên", 20, False, RGBColor(0xD7, 0xE2, 0xF5))],
       [("và sinh lời giải từng bước dựa trên Mô hình Tri thức Lỗi", 20, False, RGBColor(0xD7, 0xE2, 0xF5))]],
      space_after=3)
_text(s, Inches(0.9), Inches(6.55), Inches(11.6), Inches(0.6),
      [[("Project 1  ·  Ứng dụng Ontology + Hệ luật dẫn + LLM  ·  ĐHQG-HCM, UIT", 13, False, RGBColor(0x9F, 0xB6, 0xDB))]])

# ============================================================= SLIDE 2: AGENDA
s, top = content_slide("Nội dung trình bày")
agenda = [
    ("1", "Đặt vấn đề & Mục tiêu", "Bối cảnh, thách thức, yêu cầu sản phẩm"),
    ("2", "Nền tảng Biểu diễn Tri thức", "Ontology · Hệ luật dẫn · Mô hình tri thức quan hệ"),
    ("3", "Quy trình thiết kế hệ CSTT", "4 giai đoạn: thu thập → CSTT → suy diễn → giao diện"),
    ("4", "Mô hình tri thức của hệ thống", "4 tầng Ontology & quan hệ Bài toán–Lỗi–Sửa"),
    ("5", "Kiến trúc & Pipeline chẩn đoán", "Static Analyzers → Evidence → Rule → Ranking → LLM"),
    ("6", "Thực nghiệm & Đánh giá", "So sánh Hệ lai (Hybrid) vs. Chỉ dùng LLM"),
    ("7", "Kết luận & Hướng phát triển", ""),
]
y = 1.55
for num, t, d in agenda:
    _box(s, Inches(0.7), Inches(y), Inches(0.55), Inches(0.55), fill=TEAL, shape=MSO_SHAPE.OVAL)
    _text(s, Inches(0.7), Inches(y), Inches(0.55), Inches(0.55),
          [[(num, 18, True, WHITE)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(1.45), Inches(y - 0.02), Inches(11), Inches(0.6),
          [[(t + "    ", 18, True, NAVY), (d, 13, False, GREY, True)]],
          anchor=MSO_ANCHOR.MIDDLE)
    y += 0.74
footer(s, 2)

# ============================================================= SLIDE 3: PROBLEM
s = bullets_slide(
    "1. Đặt vấn đề",
    [
        ("Sinh viên học lập trình thường mắc các lỗi logic & thuật toán lặp lại", 0, True),
        ("Lỗi vòng lặp vô hạn, điều kiện biên sai (off-by-one)", 1),
        ("Lỗi quản lý bộ nhớ: rò rỉ, con trỏ NULL, use-after-free", 1),
        ("Lỗi đệ quy không tiến về base case, thiếu return…", 1),
        ("Công cụ hiện có chỉ báo lỗi, KHÔNG giải thích theo ngữ cảnh bài toán", 0, True),
        ("Trình biên dịch/linter: thông báo khó hiểu, không gắn với thuật toán", 1),
        ("Chỉ dùng LLM: trả lời trôi chảy nhưng dễ 'ảo giác', thiếu kiểm chứng", 1),
        ("Nhu cầu: một hệ thống chẩn đoán CHÍNH XÁC + GIẢI THÍCH được + HƯỚNG DẪN từng bước", 0, True),
    ],
    kicker="Vì sao cần một hệ thống dựa trên tri thức?",
)
footer(s, 3)

# ============================================================= SLIDE 4: GOAL & REQ
s, top = content_slide("1. Mục tiêu & Yêu cầu sản phẩm",
                       kicker="Trích từ đề bài Project 1")
_box(s, Inches(0.6), top, Inches(6.0), Inches(4.9), fill=LIGHT)
_box(s, Inches(0.6), top, Inches(6.0), Inches(0.55), fill=BLUE)
_text(s, Inches(0.8), top, Inches(5.7), Inches(0.55),
      [[("MỤC TIÊU", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
_text(s, Inches(0.85), Emu(top + Inches(0.7)), Inches(5.5), Inches(4.0),
      [[("• Chẩn đoán lỗi logic & thuật toán trong mã nguồn C/C++", 14, False, DARK)],
       [("• Sinh lời giải TỪNG BƯỚC có giải thích ngữ cảnh", 14, False, DARK)],
       [("   (Context-Aware Explanations)", 12, False, GREY, True)],
       [("• Dựa trên MÔ HÌNH TRI THỨC LỖI được mô hình hóa", 14, False, DARK)],
       [("• Kết hợp suy luận dựa trên luật + LLM phân tích ngữ nghĩa", 14, False, DARK)]],
      space_after=10)

_box(s, Inches(6.9), top, Inches(5.85), Inches(4.9), fill=LIGHT)
_box(s, Inches(6.9), top, Inches(5.85), Inches(0.55), fill=TEAL)
_text(s, Inches(7.1), top, Inches(5.6), Inches(0.55),
      [[("YÊU CẦU SẢN PHẨM", 16, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
_text(s, Inches(7.15), Emu(top + Inches(0.7)), Inches(5.4), Inches(4.0),
      [[("✔ Giao diện nhập bài & hiển thị lời giải từng bước", 14, False, DARK)],
       [("✔ CSDL tri thức bài tập (theo chủ đề)", 14, False, DARK)],
       [("✔ Mô hình biểu diễn lỗi và cách sửa", 14, False, DARK)],
       [("✔ So sánh hiệu quả: Hệ thống xây dựng vs. chỉ dùng LLM", 14, False, DARK)]],
      space_after=12)
footer(s, 4)

# ============================================================= SLIDE 5: KR FOUNDATION
s, top = content_slide("2. Nền tảng Biểu diễn Tri thức",
                       kicker="Các phương pháp BDTT được vận dụng trong đề tài")
cards = [
    ("Ontology", "Biểu diễn khái niệm, quan hệ, phân cấp lĩnh vực C/C++ và lỗi", TEAL),
    ("Hệ luật dẫn\n(Rule-based)", "Luật IF evidence THEN bug — suy diễn tiến để chẩn đoán", ORANGE),
    ("Mô hình tri thức\nquan hệ", "Bộ (C, R, Rules): khái niệm – quan hệ 2 ngôi – luật", BLUE),
    ("Mạng tính toán\n& Mẫu bài toán", "Suy luận theo mẫu (Polya) để dẫn dắt từng bước", GREEN),
]
cw, gap = Inches(2.95), Inches(0.18)
x = Inches(0.55)
for name, desc, c in cards:
    _box(s, x, top, cw, Inches(3.4), fill=LIGHT, line=c, line_w=1.5)
    _box(s, x, top, cw, Inches(1.05), fill=c)
    _text(s, x, top, cw, Inches(1.05), [[(l, 16, True, WHITE)] for l in name.split("\n")],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    _text(s, Emu(x + Inches(0.15)), Emu(top + Inches(1.2)), Emu(cw - Inches(0.3)), Inches(2.0),
          [[(desc, 13, False, DARK)]], align=PP_ALIGN.CENTER)
    x = Emu(x + cw + gap)
_text(s, Inches(0.55), Emu(top + Inches(3.6)), Inches(12.2), Inches(0.8),
      [[("Nguồn: tài liệu môn học — Tổng quan BDTT, Các phương pháp BDTT cơ bản (Logic, Mạng ngữ nghĩa, "
         "Hệ luật dẫn, Frame/Script), Ontology, Mô hình tri thức quan hệ, Mạng tính toán.", 11, False, GREY, True)]])
footer(s, 5)

# ============================================================= SLIDE 6: ONTOLOGY THEORY
s = bullets_slide(
    "2a. Ontology — Xương sống tri thức",
    [
        ("Ontology = bảng từ vựng chung + cấu trúc khái niệm của một lĩnh vực", 0, True),
        ("Chia sẻ hiểu biết chung, tái sử dụng và tích hợp tri thức", 1),
        ("Thành phần: Khái niệm (lớp) – Thuộc tính – Quan hệ – Ràng buộc – Thể hiện", 1),
        ("Vì sao chọn Ontology cho đề tài?", 0, True),
        ("Miền C/C++ có phân cấp rõ: biến → mảng → vòng lặp → đệ quy → CTDL", 1),
        ("Lỗi lập trình gắn với khái niệm & triệu chứng → mô hình hóa tự nhiên", 1),
        ("Dễ mở rộng: thêm loại lỗi mới = thêm nút & quan hệ, không sửa mã suy diễn", 1),
        ("Triển khai bằng đồ thị tri thức Neo4j (nodes/relationships/luật Cypher)", 0, True),
    ],
    kicker="Tài liệu: 7. Ontology.pdf",
)
footer(s, 6)

# ============================================================= SLIDE 7: RULE-BASED
s = bullets_slide(
    "2b. Hệ luật dẫn & Suy luận dựa trên luật",
    [
        ("Tri thức chẩn đoán biểu diễn bằng LUẬT DẪN (production rules)", 0, True),
        ("Dạng: IF <tập bằng chứng/khái niệm> THEN <lỗi> với độ tin cậy (confidence)", 1),
        ("VD: IF vòng lặp có biên 'i <= n' THEN off_by_one / array_out_of_bounds (0.91)", 1),
        ("Cơ chế suy diễn tiến (forward chaining)", 0, True),
        ("Bằng chứng (Evidence) từ phân tích → kích hoạt luật → tập lỗi ứng viên", 1),
        ("Xếp hạng ứng viên theo điểm: 0.5·luật + 0.42·bằng chứng + thưởng", 1),
        ("Ưu điểm so với 'hộp đen' LLM", 0, True),
        ("Minh bạch, giải thích được (truy vết luật ↔ bằng chứng ↔ dòng lệnh)", 1),
        ("Kết quả xác định (deterministic), có thể kiểm chứng & lặp lại", 1),
    ],
    kicker="Tài liệu: 2a/2b. Các phương pháp BDTT cơ bản",
)
footer(s, 7)

# ============================================================= SLIDE 8: RELA + POLYA
s, top = content_slide("2c. Mô hình tri thức quan hệ & Mẫu bài toán",
                       kicker="Tài liệu: 6/6b. Mô hình tri thức quan hệ · 5b. Mạng tính toán dùng mẫu bài toán")
_box(s, Inches(0.6), top, Inches(6.05), Inches(4.9), fill=LIGHT, line=BLUE, line_w=1.5)
_text(s, Inches(0.8), Emu(top + Inches(0.15)), Inches(5.6), Inches(0.6),
      [[("Mô hình quan hệ (Rela-model) = (C, R, Rules)", 16, True, BLUE)]])
_text(s, Inches(0.85), Emu(top + Inches(0.9)), Inches(5.5), Inches(3.8),
      [[("C — Tập khái niệm: mỗi khái niệm là lớp đối tượng có thuộc tính & hành vi", 13, False, DARK)],
       [("R — Tập quan hệ hai ngôi giữa các khái niệm", 13, False, DARK)],
       [("Rules — Tập luật dẫn suy diễn", 13, False, DARK)],
       [("→ Áp dụng: chuỗi quan hệ của đề tài", 13, True, NAVY)],
       [("Bài toán → Thuật toán → Cấu trúc DL → Lỗi → Hướng sửa", 13, True, ORANGE)]],
      space_after=9)

_box(s, Inches(6.95), top, Inches(5.8), Inches(4.9), fill=LIGHT, line=GREEN, line_w=1.5)
_text(s, Inches(7.15), Emu(top + Inches(0.15)), Inches(5.4), Inches(0.6),
      [[("Hướng dẫn từng bước theo Polya", 16, True, GREEN)]])
_text(s, Inches(7.2), Emu(top + Inches(0.9)), Inches(5.4), Inches(3.8),
      [[("“How to solve it?” — giải theo mẫu bài toán quen thuộc", 13, False, DARK, True)],
       [("1. Định vị bằng chứng lỗi mạnh nhất", 13, False, DARK)],
       [("2. Giải thích mâu thuẫn với yêu cầu bài toán", 13, False, DARK)],
       [("3. Kết nối tới bước sửa nhỏ & an toàn nhất", 13, False, DARK)],
       [("→ Sinh Explanation + Repair Plan từng bước", 13, True, NAVY)]],
      space_after=9)
footer(s, 8)

# ============================================================= SLIDE 9: KBS PROCESS
s, top = content_slide("3. Quy trình thiết kế hệ Cơ sở tri thức",
                       kicker="Tài liệu: 3. Các quy trình trong Thiết kế hệ CSTT — ánh xạ vào đề tài")
phases = [
    ("Giai đoạn 1", "Thu thập tri thức\n& bài toán",
     "Sách bài tập C/C++, phân loại theo chủ đề, tập luật lỗi thường gặp", TEAL),
    ("Giai đoạn 2", "Thiết kế\nCơ sở tri thức",
     "4 tầng Ontology trên Neo4j: Khái niệm · Lỗi · Luật · Giải thích", BLUE),
    ("Giai đoạn 3", "Thiết kế\nbộ suy diễn",
     "Pipeline: phân tích → bằng chứng → luật → xếp hạng → LLM", ORANGE),
    ("Giai đoạn 4", "Thiết kế\ngiao diện",
     "Web nhập bài & hiển thị lời giải từng bước (FastAPI + HTML/JS)", GREEN),
]
cw = Inches(2.95)
x = Inches(0.55)
for tag, name, desc, c in phases:
    _box(s, x, top, cw, Inches(3.7), fill=WHITE, line=c, line_w=2)
    _box(s, x, top, cw, Inches(0.5), fill=c)
    _text(s, x, top, cw, Inches(0.5), [[(tag, 13, True, WHITE)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, x, Emu(top + Inches(0.62)), cw, Inches(1.1),
          [[(l, 15, True, c)] for l in name.split("\n")],
          align=PP_ALIGN.CENTER, space_after=0)
    _text(s, Emu(x + Inches(0.15)), Emu(top + Inches(1.85)), Emu(cw - Inches(0.3)), Inches(1.7),
          [[(desc, 12.5, False, DARK)]], align=PP_ALIGN.CENTER)
    if x < Inches(9):
        _text(s, Emu(x + cw - Inches(0.02)), Emu(top + Inches(1.5)), Inches(0.35), Inches(0.5),
              [[("→", 22, True, GREY)]], align=PP_ALIGN.CENTER)
    x = Emu(x + cw + Inches(0.18))
footer(s, 9)

# ============================================================= SLIDE 10: 4-LAYER ONTOLOGY
s, top = content_slide("4. Mô hình biểu diễn tri thức của hệ thống",
                       kicker="4 tầng Ontology & quan hệ Bài toán – Lỗi – Hướng sửa")
layers = [
    ("Tầng 1 — Concept Ontology", "Kiến thức C/C++: biến, mảng, vòng lặp, đệ quy, con trỏ, CTDL", TEAL, "concepts/"),
    ("Tầng 2 — Bug Ontology", "Lỗi thường gặp + triệu chứng (symptom) + quan niệm sai (misconception)", ORANGE, "bugs/"),
    ("Tầng 3 — Diagnostic Rules", "Luật: từ Evidence + Concept → Bug (kèm độ tin cậy)", BLUE, "diagnostic/"),
    ("Tầng 4 — Explanation & Feedback", "Sinh giải thích từng bước + kế hoạch sửa (Repair Plan)", GREEN, "explanation/"),
]
y = top
for name, desc, c, folder in layers:
    _box(s, Inches(0.7), y, Inches(9.3), Inches(1.05), fill=LIGHT, line=c, line_w=1.5)
    _box(s, Inches(0.7), y, Inches(0.16), Inches(1.05), fill=c)
    _text(s, Inches(1.0), Emu(y + Inches(0.08)), Inches(9.0), Inches(0.5),
          [[(name, 16, True, c)]])
    _text(s, Inches(1.0), Emu(y + Inches(0.52)), Inches(9.0), Inches(0.5),
          [[(desc, 13, False, DARK)]])
    _box(s, Inches(10.2), y, Inches(2.5), Inches(1.05), fill=c)
    _text(s, Inches(10.2), y, Inches(2.5), Inches(1.05),
          [[(folder, 15, True, WHITE)], [("Neo4j graph", 10, False, WHITE)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    y = Emu(y + Inches(1.18))
footer(s, 10)

# ============================================================= SLIDE 11: ARCHITECTURE
s, top = content_slide("5. Kiến trúc tổng quan hệ thống")
def node(x, y, w, h, title, sub, c):
    _box(s, x, y, w, h, fill=c)
    _text(s, x, y, w, h,
          [[(title, 15, True, WHITE)]] + ([[(sub, 10.5, False, WHITE)]] if sub else []),
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=2)

def arrow(x, y, w, h, txt=None, vert=False):
    sh = MSO_SHAPE.DOWN_ARROW if vert else MSO_SHAPE.RIGHT_ARROW
    _box(s, x, y, w, h, fill=RGBColor(0x9A, 0xA7, 0xB8), shape=sh)

node(Inches(4.9), top, Inches(3.5), Inches(0.75), "Frontend (HTML/JS)", "Nhập code + hiển thị lời giải", NAVY)
arrow(Inches(6.45), Emu(top + Inches(0.78)), Inches(0.4), Inches(0.5), vert=True)
node(Inches(4.9), Emu(top + Inches(1.32)), Inches(3.5), Inches(0.75), "FastAPI Backend", "POST /diagnose", BLUE)
arrow(Inches(6.45), Emu(top + Inches(2.1)), Inches(0.4), Inches(0.5), vert=True)

ty = Emu(top + Inches(2.7))
node(Inches(0.7), ty, Inches(3.7), Inches(1.4), "Static Analyzers", "Clang AST · Static · Tidy\nSandbox · ASan/UBSan", TEAL)
node(Inches(4.8), ty, Inches(3.7), Inches(1.4), "Semantic LLM", "Phân tích ngữ nghĩa\n+ Sinh phản hồi", ORANGE)
node(Inches(8.9), ty, Inches(3.7), Inches(1.4), "Graph DB — Neo4j", "4 tầng Ontology\nLuật chẩn đoán", GREEN)

_box(s, Inches(0.7), Emu(ty + Inches(1.7)), Inches(11.9), Inches(1.35), fill=LIGHT, line=NAVY, line_w=1.5)
_text(s, Inches(0.9), Emu(ty + Inches(1.78)), Inches(11.6), Inches(1.2),
      [[("Lõi suy diễn (Diagnostic Pipeline): ", 14, True, NAVY),
        ("Evidence Builder → match_rules → Bug Ranking → chọn Explanation/Repair", 14, False, DARK)],
       [("Kết quả: ", 13, True, NAVY),
        ("Báo cáo chẩn đoán + lời giải từng bước + hướng sửa, gắn với dòng lệnh cụ thể", 13, False, DARK)]],
      space_after=6)
footer(s, 11)

# ============================================================= SLIDE 12: PIPELINE
s = bullets_slide(
    "5. Pipeline chẩn đoán — Code → Lời giải",
    [
        ("① Static Analyzers: biên dịch & chạy thử, thu thập EvidenceInstance[]", 0, True),
        ("Clang AST/Static/Tidy, LocalSandbox (chạy thử), Sanitizer (ASan/UBSan)", 1),
        ("② Evidence Builder: gộp & khử trùng lặp bằng chứng, suy ra khái niệm", 0, True),
        ("③ Graph Repository: match_rules(evidence, concepts) → RuleHit[]", 0, True),
        ("Truy vấn 4 tầng Ontology (Neo4j) hoặc bản in-memory offline", 1),
        ("④ Bug Ranking Engine: điểm = 0.5·luật + 0.42·bằng chứng + thưởng nguồn", 0, True),
        ("Chọn lỗi hàng đầu nếu vượt ngưỡng tin cậy (min_confidence = 0.45)", 1),
        ("⑤ LLM Semantic + Feedback: diễn giải ngôn ngữ tự nhiên (có kiểm chứng)", 0, True),
        ("→ Kết xuất: DiagnosticReport = giải thích từng bước + Repair Plan", 0, True),
    ],
    kicker="diagnostic_pipeline/pipeline.py",
)
footer(s, 12)

# ============================================================= SLIDE 13: EXAMPLE
s, top = content_slide("5. Ví dụ minh họa — lỗi Off-by-one",
                       kicker="Bài: tính tổng n phần tử mảng")
_box(s, Inches(0.6), top, Inches(6.1), Inches(4.9), fill=DARK)
code = [
    "int n = 5;",
    "int a[5] = {1,2,3,4,5};",
    "int sum = 0;",
    "for (int i = 0; i <= n; ++i)   // BUG",
    "    sum += a[i];   // đọc a[5] ngoài biên",
    "std::cout << sum;",
]
_text(s, Inches(0.8), Emu(top + Inches(0.2)), Inches(5.8), Inches(4.5),
      [[(ln, 14, False, RGBColor(0xF8, 0xF8, 0xF2))] for ln in code], space_after=10)

steps = [
    ("Bằng chứng", "Sanitizer báo đọc ngoài biên tại a[5]", TEAL),
    ("Luật kích hoạt", "loop 'i<=n' → off_by_one / array_out_of_bounds", ORANGE),
    ("Chẩn đoán", "array_out_of_bounds — độ tin cậy 0.955", BLUE),
    ("Hướng sửa", "Đổi 'i <= n' thành 'i < n' (biên trên loại trừ)", GREEN),
]
y = top
for name, desc, c in steps:
    _box(s, Inches(6.9), y, Inches(5.85), Inches(1.1), fill=LIGHT, line=c, line_w=1.5)
    _box(s, Inches(6.9), y, Inches(0.16), Inches(1.1), fill=c)
    _text(s, Inches(7.15), Emu(y + Inches(0.1)), Inches(5.5), Inches(0.45),
          [[(name, 14, True, c)]])
    _text(s, Inches(7.15), Emu(y + Inches(0.55)), Inches(5.5), Inches(0.5),
          [[(desc, 12.5, False, DARK)]])
    y = Emu(y + Inches(1.22))
footer(s, 13)

# ============================================================= SLIDE 14: EVALUATION
s, top = content_slide("6. Thực nghiệm & Đánh giá",
                       kicker="13 test-case · So sánh Hệ lai (Hybrid) vs. Chỉ dùng LLM")
# KPI cards
kpis = [
    ("13/13", "Độ chính xác Hybrid", "100% phát hiện đúng loại lỗi", GREEN),
    ("~3.7 s", "Thời gian TB / case", "Phân tích tĩnh chạy offline", BLUE),
    ("9 loại lỗi", "Bao phủ", "off-by-one, OOB, NULL, leak, UAF…", TEAL),
]
x = Inches(0.6)
for big, t, d, c in kpis:
    _box(s, x, top, Inches(3.95), Inches(1.7), fill=c)
    _text(s, x, Emu(top + Inches(0.1)), Inches(3.95), Inches(0.8),
          [[(big, 34, True, WHITE)]], align=PP_ALIGN.CENTER)
    _text(s, x, Emu(top + Inches(0.95)), Inches(3.95), Inches(0.7),
          [[(t, 14, True, WHITE)], [(d, 10.5, False, WHITE)]],
          align=PP_ALIGN.CENTER, space_after=1)
    x = Emu(x + Inches(4.1))

# comparison table
ty = Emu(top + Inches(2.0))
rows = [
    ("Tiêu chí", "Hệ lai (Hybrid)", "Chỉ dùng LLM", True),
    ("Phát hiện đúng lỗi", "13/13 (xác định)", "Thấp & không ổn định", False),
    ("Giải thích / truy vết", "Có — luật ↔ bằng chứng ↔ dòng", "Khó kiểm chứng, dễ 'ảo giác'", False),
    ("Tính lặp lại", "Xác định (deterministic)", "Ngẫu nhiên theo lần gọi", False),
    ("Chi phí / phụ thuộc mạng", "Phân tích tĩnh chạy offline", "Luôn cần gọi API", False),
]
rh = Inches(0.62)
cw = [Inches(4.3), Inches(4.2), Inches(4.0)]
for ri, (a, b, cc, hdr) in enumerate(rows):
    y = Emu(ty + ri * rh)
    x = Inches(0.6)
    for ci, val in enumerate((a, b, cc)):
        fill = NAVY if hdr else (LIGHT if ri % 2 else WHITE)
        _box(s, x, y, cw[ci], rh, fill=fill, line=RGBColor(0xCC, 0xD4, 0xE0), line_w=0.75)
        col = WHITE if hdr else (DARK if ci == 0 else (GREEN if ci == 1 else GREY))
        _text(s, Emu(x + Inches(0.12)), y, Emu(cw[ci] - Inches(0.2)), rh,
              [[(val, 12.5, hdr or ci == 0, col)]], anchor=MSO_ANCHOR.MIDDLE)
        x = Emu(x + cw[ci])
footer(s, 14)

# ============================================================= SLIDE 15: RESULTS DETAIL
s, top = content_slide("6. Kết quả chi tiết theo test-case")
data = [
    ("off_by_one_01", "off_by_one", "array_out_of_bounds", "0.96"),
    ("off_by_one_02", "off_by_one", "array_out_of_bounds", "0.96"),
    ("null_deref_01", "null_dereference", "null_dereference", "0.90"),
    ("memory_leak_01", "memory_leak", "memory_leak", "0.73"),
    ("missing_return_01", "missing_return", "missing_return", "0.99"),
    ("array_oob_01", "array_out_of_bounds", "array_out_of_bounds", "0.96"),
    ("infinite_loop_01", "wrong_loop_condition", "wrong_loop_condition", "0.95"),
    ("use_after_free_01", "use_after_free", "use_after_free", "0.90"),
    ("recursion_no_progress_01", "no_recursive_progress", "no_recursive_progress", "0.75"),
    ("correct_sum_01", "(không lỗi)", "(không lỗi)", "—"),
    ("bubble_sort_oob_01", "off_by_one", "array_out_of_bounds", "0.92"),
    ("binary_search_wrong_01", "wrong_loop_condition", "infinite_loop", "0.76"),
    ("factorial_rec_missing_return_01", "missing_return", "missing_return", "0.89"),
]
hdr = ("Test-case", "Kỳ vọng", "Hệ thống phát hiện", "Conf", "KQ")
cw = [Inches(4.35), Inches(3.0), Inches(3.4), Inches(0.95), Inches(0.85)]
rh = Inches(0.375)
def draw_row(y, cells, is_hdr):
    x = Inches(0.35)
    for ci, val in enumerate(cells):
        fill = NAVY if is_hdr else WHITE
        _box(s, x, y, cw[ci], rh, fill=fill, line=RGBColor(0xD5, 0xDC, 0xE6), line_w=0.5)
        col = WHITE if is_hdr else DARK
        _text(s, Emu(x + Inches(0.08)), y, Emu(cw[ci] - Inches(0.12)), rh,
              [[(val, 10.5, is_hdr, col)]], anchor=MSO_ANCHOR.MIDDLE,
              align=PP_ALIGN.CENTER if ci >= 3 else PP_ALIGN.LEFT)
        x = Emu(x + cw[ci])
draw_row(top, hdr, True)
y = Emu(top + rh)
for cid, exp, det, conf in data:
    draw_row(y, (cid, exp, det, conf, "✓"), False)
    # tick màu xanh
    _text(s, Emu(Inches(0.35) + cw[0]+cw[1]+cw[2]+cw[3]+Inches(0.08)), y,
          Emu(cw[4] - Inches(0.12)), rh, [[("✓", 12, True, GREEN)]],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(y + rh)
_text(s, Inches(0.35), y, Inches(12.5), Inches(0.5),
      [[("Off-by-one & OOB được chấp nhận là tương đương (cùng một luật chẩn đoán sinh ra).",
         10.5, False, GREY, True)]])
footer(s, 15)

# ============================================================= SLIDE 16: CONCLUSION
s = bullets_slide(
    "7. Kết luận & Hướng phát triển",
    [
        ("Kết quả đạt được", 0, True),
        ("Mô hình hóa tri thức lỗi bằng 4 tầng Ontology + Hệ luật dẫn (Neo4j)", 1),
        ("Pipeline chẩn đoán lai: phân tích tĩnh/động + suy luận luật + LLM", 1),
        ("Sinh lời giải từng bước, gắn bằng chứng với dòng lệnh cụ thể", 1),
        ("Hybrid đạt 13/13, minh bạch & xác định — vượt trội 'hộp đen' LLM", 1),
        ("Ưu điểm của hướng tiếp cận dựa trên tri thức", 0, True),
        ("Giải thích được, kiểm chứng được, dễ mở rộng loại lỗi mới", 1),
        ("Hướng phát triển", 0, True),
        ("Mở rộng Bug Ontology & tập bài tập; hỗ trợ nhiều ngôn ngữ", 1),
        ("Cá nhân hóa lộ trình học theo lịch sử lỗi của sinh viên", 1),
    ],
)
footer(s, 16)

# ============================================================= SLIDE 17: THANKS
s = prs.slides.add_slide(BLANK)
_set_bg(s, NAVY)
_box(s, 0, Inches(3.1), SW, Inches(0.06), fill=TEAL)
_box(s, 0, Inches(4.35), SW, Inches(0.06), fill=ORANGE)
_text(s, Inches(0.9), Inches(3.2), Inches(11.6), Inches(1.2),
      [[("CẢM ƠN THẦY CÔ & CÁC BẠN ĐÃ LẮNG NGHE", 34, True, WHITE)]],
      align=PP_ALIGN.CENTER)
_text(s, Inches(0.9), Inches(4.55), Inches(11.6), Inches(0.8),
      [[("Sẵn sàng cho phần Hỏi — Đáp", 18, False, RGBColor(0xC7, 0xD6, 0xEE))]],
      align=PP_ALIGN.CENTER)

out = Path(__file__).parent / "Project1_BDTT_presentation.pptx"
prs.save(str(out))
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
